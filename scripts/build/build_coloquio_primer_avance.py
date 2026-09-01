#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Genera presentacion PPTX para el Primer Avance de Tesis de Maestria
Simulacion Acelerada de Speckle Optico mediante PINNs
Roberto Hernandez Estrada - UJAT DACYTI
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Rutas ──────────────────────────────────────────────────────────────────
BASE_DIR = r"D:\Tesis_Maestria"
FIG_DIR  = os.path.join(BASE_DIR, "results", "figures")
OUT_PATH = os.path.join(BASE_DIR, "slides", "coloquio_primer_avance.pptx")

FIG_1D_RES = os.path.join(FIG_DIR, "resultados_pinn_1d.png")
FIG_1D_MET = os.path.join(FIG_DIR, "metricas_adicionales_1d.png")
FIG_2D_RES = os.path.join(FIG_DIR, "resultados_pinn_2d.png")
FIG_2D_MET = os.path.join(FIG_DIR, "metricas_adicionales_2d.png")
FIG_SP_RES = os.path.join(FIG_DIR, "resultados_speckle_nb03.png")
FIG_SP_EST = os.path.join(FIG_DIR, "estadistica_speckle_nb03.png")

# ── Paleta neutral profesional ─────────────────────────────────────────────
C_BG       = RGBColor(0xFA, 0xFA, 0xF8)   # blanco hueso
C_TITLE    = RGBColor(0x1A, 0x1A, 0x2E)   # azul marino profundo
C_BODY     = RGBColor(0x2C, 0x2C, 0x2C)   # gris oscuro
C_ACCENT   = RGBColor(0x34, 0x5C, 0x8C)   # azul acero (acento principal)
C_ACCENT2  = RGBColor(0x7A, 0x92, 0xAA)   # azul-gris claro
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY    = RGBColor(0xF0, 0xF2, 0xF5)   # gris muy claro
C_MGRAY    = RGBColor(0xCC, 0xCC, 0xCC)   # gris medio
C_BOLD_NUM = RGBColor(0x1B, 0x4F, 0x72)   # azul oscuro para cifras clave
C_GREEN    = RGBColor(0x1A, 0x5C, 0x2A)   # verde exito
C_AMBER    = RGBColor(0x99, 0x55, 0x00)   # advertencia
C_SLIDE_N  = RGBColor(0xAA, 0xAA, 0xAA)   # numero de slide

# ── Dimensiones 16:9 ───────────────────────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)

# ── Presentacion ───────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

def blank():
    return prs.slide_layouts[6]

def bg(slide, color=None):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color or C_BG

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, size=18, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name  = "Calibri"
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color or C_BODY
    return box

def header(slide, title, sub=None, bar=None):
    bc = bar or C_ACCENT
    rect(slide, 0, 0, SW, Inches(1.05), fill=bc)
    tb(slide, title, Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.78),
       size=27, bold=True, color=C_WHITE)
    if sub:
        tb(slide, sub, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.25),
           size=12, italic=True, color=RGBColor(0xBB, 0xCF, 0xE5))

def slide_num(slide, n):
    tb(slide, str(n), Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
       size=11, color=C_SLIDE_N, align=PP_ALIGN.RIGHT)

def bullets(slide, items, l, t, w, h, size=18, color=None, bchar="  "):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        p.space_after  = Pt(2)
        r = p.add_run()
        if isinstance(item, tuple):
            txt, sz, bold, col = item
            r.text = txt
            r.font.size  = Pt(sz)
            r.font.bold  = bold
            r.font.color.rgb = col or color or C_BODY
        else:
            r.text = bchar + item
            r.font.size  = Pt(size)
            r.font.color.rgb = color or C_BODY
        r.font.name = "Calibri"

def metric_box(slide, label, value, unit, l, t, w=Inches(2.2), h=Inches(1.4)):
    rect(slide, l, t, w, h, fill=C_LGRAY, line=C_ACCENT, lw=Pt(1.5))
    tb(slide, value,  l+Inches(0.08), t+Inches(0.08),
       w-Inches(0.16), Inches(0.62), size=30, bold=True,
       color=C_BOLD_NUM, align=PP_ALIGN.CENTER)
    tb(slide, unit,   l+Inches(0.08), t+Inches(0.7),
       w-Inches(0.16), Inches(0.28), size=12,
       color=C_ACCENT2, align=PP_ALIGN.CENTER)
    tb(slide, label,  l+Inches(0.08), t+Inches(0.98),
       w-Inches(0.16), Inches(0.32), size=12, bold=True,
       color=C_BODY, align=PP_ALIGN.CENTER)

def figure(slide, path, l, t, w, h, cap=None):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, width=w, height=h)
    else:
        rect(slide, l, t, w, h, fill=C_LGRAY, line=C_MGRAY, lw=Pt(1))
        tb(slide, "[" + os.path.basename(path) + "]",
           l+Inches(0.1), t+h/2-Inches(0.2), w-Inches(0.2), Inches(0.4),
           size=11, color=C_MGRAY, align=PP_ALIGN.CENTER)
    if cap:
        tb(slide, cap, l, t+h+Inches(0.02), w, Inches(0.28),
           size=10, italic=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 1 — PORTADA
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s, C_ACCENT)
rect(s, 0, 0, SW, Inches(0.12), fill=C_WHITE)
rect(s, 0, SH-Inches(0.12), SW, Inches(0.12), fill=C_WHITE)

tb(s, "UNIVERSIDAD JUAREZ AUTONOMA DE TABASCO",
   Inches(0.8), Inches(0.22), Inches(11.7), Inches(0.45),
   size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Division Academica de Ciencias y Tecnologias de la Informacion (DACYTI)",
   Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.38),
   size=12, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(1.5), Inches(1.18), Inches(10.3), Inches(0.025),
     fill=RGBColor(0xBB, 0xCC, 0xDD))

tb(s, "Simulacion Acelerada de Speckle Optico",
   Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.75),
   size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "mediante Redes Neuronales Informadas por Fisica con Activacion Sinusoidal",
   Inches(0.5), Inches(2.08), Inches(12.3), Inches(0.6),
   size=20, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(2.0), Inches(2.82), Inches(9.3), Inches(0.03),
     fill=RGBColor(0xCC, 0xDD, 0xEE))

tb(s, "PRIMER AVANCE  -  COLOQUIO DE MAESTRIA",
   Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.38),
   size=14, bold=True, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

tb(s, "Roberto Hernandez Estrada",
   Inches(0.5), Inches(3.48), Inches(12.3), Inches(0.48),
   size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Maestria en Ciencias de la Computacion  |  Matricula: 252H21004",
   Inches(0.5), Inches(3.94), Inches(12.3), Inches(0.35),
   size=13, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(2.0), Inches(4.40), Inches(9.3), Inches(0.025),
     fill=RGBColor(0x80, 0xA0, 0xBB))

tb(s, "Director: Dr. Jose Adan Hernandez Nolasco",
   Inches(0.5), Inches(4.52), Inches(12.3), Inches(0.35),
   size=15, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Codirector: Dr. Oscar Alberto Chavez Bosquez",
   Inches(0.5), Inches(4.86), Inches(12.3), Inches(0.35),
   size=15, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Cunduacan, Tabasco  -  Junio 2026",
   Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.38),
   size=13, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

notes(s, """PORTADA - ~1 min
Buenos dias/tardes. Mi nombre es Roberto Hernandez Estrada, alumno del programa de Maestria en Ciencias de la Computacion de la DACYTI-UJAT.

Hoy presento el primer avance de mi tesis titulada: Simulacion Acelerada de Speckle Optico mediante Redes Neuronales Informadas por Fisica con Activacion Sinusoidal, bajo la direccion del Dr. Jose Adan Hernandez Nolasco y en codireccion del Dr. Oscar Alberto Chavez Bosquez.

En los proximos ~20 minutos les presento la motivacion, el enfoque metodologico y los resultados computacionales obtenidos hasta este punto.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Contenido de la presentacion")
slide_num(s, 2)

agenda = [
    "1.  Motivacion y el problema computacional",
    "2.  Objetivos del avance",
    "3.  Marco teorico: Speckle optico y ecuacion de Helmholtz",
    "4.  PINNs y arquitectura SIREN",
    "5.  Metodologia: muestreo, perdida, optimizacion",
    "6.  Resultados NB01 - Validacion Helmholtz 1D",
    "7.  Resultados NB02 - Helmholtz 2D campo complejo",
    "8.  Resultados NB03 - Simulacion de speckle optico",
    "9.  Comparativa con el estado del arte",
    "10. Trabajo futuro y conclusiones preliminares",
]
box = s.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(6.0))
tf  = box.text_frame
tf.word_wrap = True
for i, item in enumerate(agenda):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(9)
    r = p.add_run()
    r.text = item
    r.font.name = "Calibri"
    r.font.size = Pt(20)
    r.font.color.rgb = C_BODY

notes(s, """AGENDA - 0.5 min
Esta es la estructura de la presentacion. Comenzamos con la motivacion y el marco teorico necesario, y la mayor parte del tiempo la dedicamos a los tres experimentos y sus resultados. Terminamos con la comparativa y los proximos pasos.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 3 — MOTIVACION
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Motivacion: Por que simular speckle optico?")
slide_num(s, 3)

# Panel izq: aplicaciones
tb(s, "El speckle optico aparece en:",
   Inches(0.4), Inches(1.15), Inches(5.8), Inches(0.38),
   size=18, bold=True, color=C_ACCENT)
apps = [
    "Metrologia optica y deteccion de defectos superficiales",
    "Tomografia de coherencia optica (OCT)",
    "Imagenologia biomedica",
    "Criptografia optica",
    "Comunicaciones por fibra optica",
]
bullets(s, apps, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.2), size=18)

# Panel der: problema
rect(s, Inches(6.8), Inches(1.1), Inches(6.1), Inches(2.7),
     fill=C_LGRAY, line=C_ACCENT, lw=Pt(1))
tb(s, "El reto computacional con metodos clasicos (FEM):",
   Inches(7.0), Inches(1.2), Inches(5.7), Inches(0.38),
   size=17, bold=True, color=C_ACCENT)
probs = [
    "Malla debe tener resolucion < longitud de onda",
    "Millones de incognitas por realizacion",
    "Cada nuevo perfil de rugosidad: resolver desde cero",
    "Para estudios estadisticos (miles de realizaciones): costo intratable",
]
bullets(s, probs, Inches(7.0), Inches(1.65), Inches(5.7), Inches(2.0), size=16)

# Propuesta
rect(s, Inches(6.8), Inches(3.95), Inches(6.1), Inches(2.0),
     fill=C_ACCENT, line=None)
tb(s, "Propuesta: PINN libre de malla",
   Inches(7.0), Inches(4.05), Inches(5.7), Inches(0.38),
   size=17, bold=True, color=C_WHITE)
tb(s, "Incorporar la fisica directamente en la funcion de perdida "
      "de una red neuronal. Una vez entrenada, evalua el campo "
      "completo en ~1.3 ms (GPU RTX 5050).",
   Inches(7.0), Inches(4.45), Inches(5.7), Inches(1.35),
   size=16, color=C_WHITE)

notes(s, """MOTIVACION - 1.5 min
El speckle optico es el patron granular que emerge cuando luz coherente -como la de un laser- se dispersa sobre una superficie rugosa.

Esta presente en aplicaciones criticas: metrologia, imagen biomedica, comunicaciones opticas.

El problema computacional: simular el speckle con precision requiere resolver la ecuacion de Helmholtz con condiciones de frontera estocasticas. El Metodo de Elementos Finitos (FEM) es preciso pero muy costoso: la malla necesita resolucion menor a la longitud de onda del laser, generando millones de incognitas. Y cada nuevo perfil de rugosidad requiere resolver desde cero.

Para estudios estadisticos que requieren miles de realizaciones, esto es intratable.

Nuestra propuesta: usar una red neuronal que aprende la solucion de la EDP sin ninguna malla, evaluable en milisegundos una vez entrenada.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 4 — OBJETIVOS
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Objetivos de investigacion")
slide_num(s, 4)

tb(s, "Objetivo general",
   Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.38),
   size=20, bold=True, color=C_ACCENT)
rect(s, Inches(0.4), Inches(1.58), Inches(12.5), Inches(0.9),
     fill=C_LGRAY, line=C_ACCENT2, lw=Pt(1))
tb(s, "Desarrollar y validar una PINN-SIREN que simule speckle optico resolviendo la "
      "ecuación de Helmholtz 2D con error L² < 5% y contraste C ≈ 1.",
   Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.75),
   size=18, color=C_BODY)

rect(s, Inches(0.4), Inches(2.6), Inches(12.5), Inches(0.025), fill=C_MGRAY)

tb(s, "Objetivos especificos de este avance",
   Inches(0.4), Inches(2.7), Inches(12.5), Inches(0.38),
   size=20, bold=True, color=C_ACCENT)

objs = [
    ("NB01", "Validar la arquitectura SIREN con solucion analitica exacta (Helmholtz 1D)", C_ACCENT),
    ("NB02", "Extender al caso 2D con campo electrico complejo y muestreo LHS", C_ACCENT),
    ("NB03", "Generar patrones de speckle con validacion estadistica (criterio de Goodman)", C_ACCENT),
    ("Meta", "Error L2 < 5%  |  Contraste |C - 1| < 0.1", C_BOLD_NUM),
]
top = 3.18
for code, desc, col in objs:
    rect(s, Inches(0.4), Inches(top), Inches(0.95), Inches(0.57),
         fill=col, line=None)
    tb(s, code, Inches(0.4), Inches(top+0.03),
       Inches(0.95), Inches(0.5), size=14, bold=True,
       color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(s, desc, Inches(1.5), Inches(top+0.08),
       Inches(11.3), Inches(0.42), size=18, color=C_BODY)
    top += 0.75

notes(s, """OBJETIVOS - 1 min
El objetivo general es demostrar que una PINN-SIREN puede simular speckle optico con precision cuantitativa: error L2 menor al 5% para Helmholtz, y contraste del speckle cercano a 1 segun el criterio de Goodman.

En este primer avance cubrimos tres experimentos de complejidad creciente. Les mostrare que los tres objetivos se cumplieron.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 5 — HELMHOLTZ Y SPECKLE
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Marco teorico: Ecuacion de Helmholtz y speckle optico")
slide_num(s, 5)

# Izq: Helmholtz
tb(s, "Ecuacion de Helmholtz escalar 2D:",
   Inches(0.4), Inches(1.15), Inches(6.0), Inches(0.38), size=18, bold=True, color=C_ACCENT)
rect(s, Inches(0.4), Inches(1.58), Inches(6.0), Inches(1.6),
     fill=C_LGRAY, line=C_ACCENT2, lw=Pt(1))
tb(s, "∇²E(r) + k²E(r) = 0",
   Inches(0.6), Inches(1.72), Inches(5.6), Inches(0.5),
   size=24, bold=True, color=C_BOLD_NUM, align=PP_ALIGN.CENTER)
tb(s, "k = 2π/λ,   λ = 638 nm (láser diodo rojo)\n"
      "Dominio: Ω = [0,1]²  (adimensional)\n"
      "Campo complejo: E = E_real + i·E_imag",
   Inches(0.6), Inches(2.2), Inches(5.6), Inches(0.85), size=14, color=C_BODY)

tb(s, "La red predice ambas componentes independientemente:",
   Inches(0.4), Inches(3.3), Inches(6.0), Inches(0.35), size=16, color=C_BODY)
rect(s, Inches(0.4), Inches(3.7), Inches(6.0), Inches(0.85),
     fill=C_LGRAY, line=C_ACCENT2, lw=Pt(1))
tb(s, "∇²E_real + k²E_real = 0\n"
      "∇²E_imag + k²E_imag = 0",
   Inches(0.6), Inches(3.78), Inches(5.6), Inches(0.68),
   size=17, bold=True, color=C_BOLD_NUM)

# Der: Speckle
tb(s, "Estadistica del speckle completamente desarrollado:",
   Inches(7.0), Inches(1.15), Inches(6.0), Inches(0.38), size=18, bold=True, color=C_ACCENT)

hechos = [
    "Intensidad: I = |E|² = E²_real + E²_imag",
    "Distribución teórica: exponencial negativa",
    "p(I) = (1/⟨I⟩)·exp(−I/⟨I⟩)",
    "Contraste: C = σ_I / ⟨I⟩ = 1  (para speckle puro)",
    "Criterio validación: |C − 1| < 0.1  (Goodman 2007)",
]
bullets(s, hechos, Inches(7.0), Inches(1.6), Inches(6.0), Inches(2.8), size=17)

rect(s, Inches(7.0), Inches(4.55), Inches(6.0), Inches(2.2),
     fill=C_ACCENT, line=None)
tb(s, "Condicion de frontera NB03 (speckle):",
   Inches(7.2), Inches(4.65), Inches(5.6), Inches(0.35),
   size=15, bold=True, color=C_WHITE)
tb(s, "E(x,0) = e^(iφ(x)),   φ(x) ∼ U(0, 2π)\n"
      "N_φ = 256 puntos en y=0\n"
      "Bordes x=0, x=1, y=1: libres (sin Dirichlet)",
   Inches(7.2), Inches(5.02), Inches(5.6), Inches(1.55), size=16, color=C_WHITE)

notes(s, """MARCO TEORICO - 1.5 min
La fisica del problema esta gobernada por la ecuacion de Helmholtz escalar. Describe la propagacion de ondas electromagneticas monocromaticas en un medio homogeneo.

Como las redes neuronales operan sobre numeros reales, descomponemos el campo complejo en sus partes real e imaginaria. Cada componente satisface la misma ecuacion de forma independiente.

Para el speckle, la validacion estadistica sigue el criterio de Goodman (2007): la distribucion de intensidades debe ser exponencial negativa y el contraste C = sigma_I / <I> debe ser 1.

La condicion de frontera para el speckle es una fase aleatoria en y=0, con los otros tres bordes sin restriccion Dirichlet para evitar reflexiones artificiales.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 6 — PINNs Y SIREN
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Physics-Informed Neural Networks y arquitectura SIREN")
slide_num(s, 6)

tb(s, "Funcion de perdida PINN:",
   Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.38), size=19, bold=True, color=C_ACCENT)

rect(s, Inches(0.4), Inches(1.58), Inches(12.5), Inches(0.7),
     fill=C_LGRAY, line=C_ACCENT2, lw=Pt(1))
tb(s, "ℒ(θ) = ℒ_datos(θ)  +  λ_fis · ℒ_física(θ)",
   Inches(0.6), Inches(1.68), Inches(12.1), Inches(0.5),
   size=22, bold=True, color=C_BOLD_NUM, align=PP_ALIGN.CENTER)

tb(s, "ℒ_datos: error en condiciones de frontera",
   Inches(0.6), Inches(2.38), Inches(5.5), Inches(0.35), size=16, color=C_BODY)
tb(s, "ℒ_física: residuo de Helmholtz en puntos interiores",
   Inches(6.8), Inches(2.38), Inches(6.0), Inches(0.35), size=16, color=C_BODY)

rect(s, Inches(0.4), Inches(2.82), Inches(12.5), Inches(0.025), fill=C_MGRAY)

tb(s, "¿Por qué SIREN?  Activación sinusoidal: φ(z) = sin(ω₀·z)",
   Inches(0.4), Inches(2.92), Inches(12.5), Inches(0.38),
   size=19, bold=True, color=C_ACCENT)

probs_activ = [
    "ReLU / tanh: derivadas de 2.º orden inestables (Helmholtz requiere ∇²E)",
    "Sesgo espectral (spectral bias): aprenden frecuencias bajas primero, fallan para speckle",
    "SIREN: sin(ω₀z) → derivadas de cualquier orden son sinusoides → estabilidad numérica",
    "Inicialización especial (Sitzmann et al. 2020): distribución uniforme en [−1, 1] en todas las capas",
    "ω₀ = 1.0 (calibrado para dominio adimensional k = 2π, vs. ω₀ = 30 para imágenes)",
]
bullets(s, probs_activ, Inches(0.4), Inches(3.38), Inches(12.5), Inches(3.5), size=17)

notes(s, """PINNS Y SIREN - 2 min
Las PINNs, introducidas por Raissi, Perdikaris y Karniadakis en 2019, entrenan la red minimizando una funcion de perdida que combina el error en las condiciones de frontera con el residuo de la ecuacion diferencial en puntos interiores.

El reto para Helmholtz: necesitamos derivadas de segundo orden. Las activaciones estandar como ReLU o tanh producen derivadas inestables. Ademas, exhiben spectral bias: aprenden frecuencias bajas primero, lo que es fatal para el speckle que tiene detalles finos.

SIREN resuelve esto: la derivada de sin es cos, que tambien es sinusoidal. Esto garantiza que los gradientes de cualquier orden sean estables y distribuidos uniformemente.

Elegimos omega_0 = 1.0 en lugar del omega_0 = 30 recomendado para imagenes, porque nuestro dominio adimensional con k = 2pi no requiere frecuencias ultra-altas, y valores mayores desestabilizan el entrenamiento con L-BFGS.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 7 — ARQUITECTURA Y METODOLOGIA
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Arquitectura SIREN 5x128 y estrategia de entrenamiento")
slide_num(s, 7)

# Izq: Arquitectura
tb(s, "Arquitectura SIREN 5 x 128",
   Inches(0.4), Inches(1.15), Inches(5.8), Inches(0.38),
   size=19, bold=True, color=C_ACCENT)
arch = [
    "Entrada:  (x, y)  →  2 valores",
    "5 capas ocultas x 128 neuronas",
    "Activación: sin(ω₀·z),  ω₀ = 1.0",
    "Salida lineal: (E_real, E_imag)  →  2 neuronas",
    "Parametros entrenables: 66,690",
    "Inicializacion: Sitzmann et al. (2020)",
]
bullets(s, arch, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.2), size=18)

# Diagrama simplificado
for i, (lbl, col) in enumerate([
    ("(x,y)", C_ACCENT2),
    ("128\nsin(ω₀z)", C_ACCENT),
    ("128\nsin(ω₀z) ×3", C_ACCENT),
    ("128\nsin(ω₀z)", C_ACCENT),
    ("(Er,Ei)", C_BOLD_NUM),
]):
    x = Inches(0.45 + i*1.12)
    rect(s, x, Inches(5.2), Inches(1.0), Inches(0.7), fill=col, line=None)
    tb(s, lbl, x+Inches(0.05), Inches(5.26), Inches(0.9), Inches(0.58),
       size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
for xi in [1.6, 2.72, 3.84, 4.96]:
    tb(s, "→", Inches(xi), Inches(5.4), Inches(0.25), Inches(0.3),
       size=14, bold=True, color=C_MGRAY, align=PP_ALIGN.CENTER)

# Der: Metodologia
tb(s, "Muestreo y optimizacion",
   Inches(6.7), Inches(1.15), Inches(6.2), Inches(0.38),
   size=19, bold=True, color=C_ACCENT)
meth = [
    "Puntos interiores: LHS 3,000 pts (NB02/NB03)",
    "  Cobertura uniforme sin patrones redundantes",
    "Puntos de frontera: 300 pts por borde (x 4 bordes)",
    "Peso de física: λ_fis = 0.1",
    "  [Ablación: 0.01 → 0.1 (óptimo) → 1.0 (colapso)]",
    "",
    "Fase 1 - Adam:",
    "  η = 10⁻³,  máx 15,000 épocas,  early stopping",
    "Fase 2 - L-BFGS (strong Wolfe):",
    "  Max 1,000 iter,  historial 100",
    "  Ajuste fino de alta precision con curvatura",
]
bullets(s, meth, Inches(6.7), Inches(1.6), Inches(6.2), Inches(4.0), size=16)

notes(s, """ARQUITECTURA Y METODOLOGIA - 2 min
La arquitectura SIREN 5x128 tiene 5 capas ocultas de 128 neuronas, con activacion sinusoidal. La entrada son las coordenadas espaciales (x,y) y la salida son las dos componentes del campo. En total, 66,690 parametros.

Para el muestreo, usamos Muestreo por Hipercubo Latino (LHS). A diferencia de una malla cartesiana, el LHS garantiza cobertura uniforme sin patrones redundantes.

La estrategia bifasica de optimizacion es critica:
- Adam explora el espacio de parametros globalmente, evitando minimos locales
- L-BFGS aprovecha la informacion de curvatura de segundo orden para un ajuste fino de alta precision

El peso de fisica lambda_fis = 0.1 fue calibrado experimentalmente. Con lambda = 1.0, los dos residuos (real e imaginario) duplican la contribucion fisica y causan explosion de gradientes en float32. Con lambda = 0.01 la fisica queda subponderada.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 8 — RESULTADOS NB01 (figura principal)
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB01: Validacion Helmholtz 1D - Comparacion con solucion analitica",
       sub="SIREN 5×64  |  k = 2π  |  E_exacta(x) = cos(kx)  |  1,000 puntos de evaluación")
slide_num(s, 8)

figure(s, FIG_1D_RES, Inches(0.25), Inches(1.12), Inches(8.55), Inches(5.65),
       cap="Fig. 4.1: PINN vs solucion analitica (arriba) - Error puntual (izq) - Curvas de perdida Adam+L-BFGS (der)")

# Metricas
tb(s, "Metricas clave", Inches(9.0), Inches(1.15), Inches(4.0), Inches(0.38),
   size=18, bold=True, color=C_ACCENT)

metric_box(s, "Error L2", "0.006%", "meta: < 5%",  Inches(9.0),  Inches(1.58))
metric_box(s, "R cuadrado", "1.000000", "ajuste perfecto", Inches(11.35), Inches(1.58))
metric_box(s, "L-BFGS", "202 iter", "de 500 max",  Inches(9.0),  Inches(3.12))
metric_box(s, "Tiempo", "251 s", "total",         Inches(11.35), Inches(3.12))

rect(s, Inches(9.0), Inches(4.65), Inches(4.1), Inches(1.55),
     fill=RGBColor(0xE8, 0xF4, 0xEA), line=C_ACCENT2, lw=Pt(1))
tb(s, "Mejora vs. estado del arte:", Inches(9.1), Inches(4.72),
   Inches(3.9), Inches(0.35), size=14, bold=True, color=C_GREEN)
tb(s, "×415 vs Schoder & Kraxberger\n(2024): error 2.49%",
   Inches(9.1), Inches(5.05), Inches(3.9), Inches(1.0),
   size=20, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

notes(s, """NB01 - 2 min
El primer experimento valida la arquitectura SIREN frente a la solucion analitica exacta de la ecuacion de Helmholtz unidimensional.

La grafica superior muestra PINN (linea roja) vs cos(kx) (azul). Son visualmente indistinguibles.

La grafica inferior izquierda muestra el error puntual: practicamente cero en todo el dominio. Maxima desviacion: 5.77 x 10^{-5}.

La grafica inferior derecha muestra las curvas de perdida: la caida abrupta alrededor de la epoca 15,000 es la transicion de Adam a L-BFGS, que refina rapidamente en solo 202 iteraciones.

Resultados: error L2 = 0.006%, casi 3 ordenes de magnitud por debajo del umbral del 5%. R^2 = 1.000000. Tiempo: 251 segundos.

Factor de mejora de 415 veces respecto a Schoder & Kraxberger (2024) con PINNs convencionales.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 9 — NB01 metricas adicionales
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB01: Analisis detallado de error y convergencia")
slide_num(s, 9)

figure(s, FIG_1D_MET, Inches(0.25), Inches(1.12), Inches(12.8), Inches(5.3),
       cap="Fig. 4.2: Histograma de error absoluto (izq) - Error relativo puntual (centro) - Correlacion PINN vs Exacta: R2 = 1.000000, Pearson = 1.000000 (der)")

tb(s, "Error maximo = 5.77e-5   |   MAE = 3.50e-5   |   "
      "Adam: 15,000 epocas   |   L-BFGS: 202 iteraciones",
   Inches(0.5), Inches(6.52), Inches(12.3), Inches(0.35),
   size=14, italic=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

notes(s, """NB01 METRICAS - 1 min
Este slide muestra el analisis detallado del error.

El histograma confirma que los errores se concentran en valores muy pequenos, del orden de 10^{-5}.

El panel central muestra que el error relativo puntual tiene una media del 0.012%, practicamente cero en todo el dominio.

La grafica PINN vs Exacta muestra correlacion lineal perfecta: R2 y Pearson = 1.000000 exactamente.

Esto nos confirma que la arquitectura y el protocolo de entrenamiento estan correctamente implementados. Podemos proceder con confianza al caso 2D.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 10 — NB02 figura principal
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB02: Helmholtz 2D con campo complejo",
       sub="SIREN 5×128  |  k = 2π  |  LHS 3,000 pts  |  E_exacta = e^(i(kₓx + k_yy))")
slide_num(s, 10)

figure(s, FIG_2D_RES, Inches(0.2), Inches(1.12), Inches(8.7), Inches(5.7),
       cap="Fig. 4.3: E_real y E_imag — solución exacta (fila sup.) vs predicción PINN (fila inf.) — Intensidad I=|E|², puntos LHS y curvas de pérdida")

# Tabla metricas NB02
tb(s, "Metricas NB02", Inches(9.1), Inches(1.15), Inches(4.0), Inches(0.38),
   size=18, bold=True, color=C_ACCENT)

rows = [
    ("Metrica",         "E_real",   "E_imag"),
    ("Error L2 (%)",    "0.214",    "0.127"),
    ("R cuadrado",      "0.999995", "0.999998"),
    ("Pearson",         "0.999998", "0.999999"),
    ("Error L2 prom.", "0.171%",   ""),
    ("Epocas Adam",     "8,737",    ""),
    ("Iter. L-BFGS",    "1,035",    ""),
    ("Tiempo total",    "299 s",    ""),
]
t_top = 1.60
for i, (m, v1, v2) in enumerate(rows):
    is_hdr = (i == 0)
    is_prom = (i == 4)
    bg_r = C_ACCENT if is_hdr else (RGBColor(0xE8, 0xF4, 0xEA) if is_prom else (C_LGRAY if i % 2 == 0 else C_WHITE))
    tc   = C_WHITE if is_hdr else (C_GREEN if is_prom else C_BODY)
    rect(s, Inches(9.1), Inches(t_top+i*0.5), Inches(4.0), Inches(0.48),
         fill=bg_r, line=C_ACCENT2, lw=Pt(0.5))
    tb(s, m,  Inches(9.15), Inches(t_top+i*0.5+0.07),
       Inches(1.7), Inches(0.35), size=12, bold=(is_hdr or is_prom), color=tc)
    tb(s, v1, Inches(10.9), Inches(t_top+i*0.5+0.07),
       Inches(1.1), Inches(0.35), size=12, bold=(is_hdr or is_prom),
       color=C_BOLD_NUM if is_prom else tc, align=PP_ALIGN.CENTER)
    tb(s, v2, Inches(12.05), Inches(t_top+i*0.5+0.07),
       Inches(0.95), Inches(0.35), size=12, bold=is_hdr,
       color=tc, align=PP_ALIGN.CENTER)

notes(s, """NB02 - 2 min
El segundo experimento escala al caso bidimensional con campo complejo.

Los mapas de calor muestran E_real y E_imag. La fila superior es la solucion analitica exacta -una onda plana diagonal- y la fila inferior es la prediccion PINN. La diferencia visual es indistinguible.

Los errores puntuales son del orden de 10^{-3}: en la tercera cifra decimal.

Resultados cuantitativos: error L2 promedio de 0.171%, por debajo del umbral del 5% por mas de un orden de magnitud. R^2 > 0.9999 para ambas componentes. Tiempo: 299 segundos.

La mejora respecto al estado del arte es de un factor de 11.2 para el caso 2D/3D.

El aumento de capacidad de 64 a 128 neuronas fue determinante: con d=64 el error era 0.436%; con d=128 bajo a 0.171%, una mejora de 2.5x.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 11 — NB02 metricas y robustez
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB02: Analisis de error y reproducibilidad multi-semilla")
slide_num(s, 11)

figure(s, FIG_2D_MET, Inches(0.2), Inches(1.12), Inches(7.6), Inches(5.15),
       cap="Fig. 4.4: Distribucion de errores E_real / E_imag - Mapa de error puntual - Correlacion PINN vs Exacta")

tb(s, "Robustez multi-semilla (NB02, λ_fis = 0.1)",
   Inches(8.0), Inches(1.15), Inches(5.0), Inches(0.38),
   size=16, bold=True, color=C_ACCENT)

seeds = [
    ("SEED", "Error L2", "Epocas", "Tiempo"),
    ("42",   "0.171%",   "8,737",  "299 s"),
    ("123",  "0.126%",   "8,012",  "448 s"),
    ("777",  "0.168%",   "13,560", "346 s"),
    ("Media ± DE", "0.155 ± 0.020%", "--", "--"),
]
t_top = 1.60
for i, row in enumerate(seeds):
    is_hdr = (i == 0)
    is_sum = (i == 4)
    bg_r = C_ACCENT if is_hdr else (RGBColor(0xE8, 0xF4, 0xEA) if is_sum else (C_LGRAY if i % 2 == 0 else C_WHITE))
    tc   = C_WHITE if is_hdr else (C_GREEN if is_sum else C_BODY)
    rect(s, Inches(8.0), Inches(t_top+i*0.55), Inches(5.0), Inches(0.52),
         fill=bg_r, line=C_ACCENT2, lw=Pt(0.5))
    col_x = [8.05, 9.3, 10.4, 11.5]
    col_w = [1.22, 1.07, 1.07, 1.42]
    for j, (val, cw) in enumerate(zip(row, col_w)):
        tb(s, val, Inches(col_x[j]), Inches(t_top+i*0.55+0.09),
           Inches(cw), Inches(0.35), size=13, bold=(is_hdr or is_sum),
           color=tc, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

rect(s, Inches(8.0), Inches(4.65), Inches(5.0), Inches(1.3),
     fill=RGBColor(0xE8, 0xF4, 0xEA), line=C_ACCENT2, lw=Pt(1))
tb(s, "Varianza inter-semilla: 13% relativo",
   Inches(8.15), Inches(4.73), Inches(4.7), Inches(0.38),
   size=16, bold=True, color=C_GREEN)
tb(s, "El metodo es estadisticamente robusto:\nconverge independientemente de la inicializacion.",
   Inches(8.15), Inches(5.1), Inches(4.7), Inches(0.75),
   size=15, color=C_BODY)

notes(s, """NB02 ROBUSTEZ - 1.5 min
Para evaluar la robustez estadistica, repetimos NB02 con tres semillas diferentes: 42, 123 y 777.

Los resultados son muy consistentes: error L2 entre 0.126% y 0.171%, con media 0.155% y desviacion estandar 0.020%.

La varianza inter-semilla del 13% relativo confirma que el metodo es estadisticamente robusto: independientemente de la inicializacion aleatoria, el modelo converge a soluciones de alta precision.

Tambien observamos en el mapa de error que los errores mas altos se concentran en las esquinas del dominio, la region geometricamente mas dificil para condiciones de onda plana.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 12 — NB03 patron speckle
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB03: Simulacion de speckle optico - Patron generado",
       sub="φ(x) ∼ U(0, 2π)  |  N_φ = 256 pts en y=0  |  Bordes x=0, x=1, y=1 libres  |  Misma arquitectura SIREN 5×128")
slide_num(s, 12)

figure(s, FIG_SP_RES, Inches(0.2), Inches(1.12), Inches(13.0), Inches(5.9),
       cap="Fig. 4.5 (izquierda): patrón de intensidad I=|E|², componentes E_real, E_imag, fase aleatoria φ(x,y), puntos LHS y curvas de pérdida")

notes(s, """NB03 PATRON - 1.5 min
El tercer experimento es el corazon de la tesis: generar un patron de speckle optico.

La unica modificacion respecto a NB02 es la condicion de frontera en y=0: en lugar de una onda plana conocida, imponemos una fase aleatoria uniforme en [0, 2pi] muestreada en 256 puntos. Esta condicion modela una superficie rugosa iluminada por laser.

Los tres bordes restantes se dejan libres para evitar reflexiones espurias.

El patron resultante muestra la textura granular caracteristica del speckle completamente desarrollado: zonas brillantes y oscuras distribuidas aleatoriamente.

Importante: la misma arquitectura SIREN 5x128 que resuelve Helmholtz con onda plana tambien resuelve la ecuacion con frontera rugosa aleatoria, sin ninguna modificacion de arquitectura ni funcion de perdida. Solo cambia la condicion de frontera. Esto demuestra la generalidad del metodo.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 13 — NB03 validacion estadistica
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB03: Validacion estadistica del speckle generado")
slide_num(s, 13)

figure(s, FIG_SP_EST, Inches(0.2), Inches(1.12), Inches(7.8), Inches(5.15),
       cap="Fig. 4.5 (derecha): distribución de intensidades normalizada vs exponencial negativa teórica e^(−I/⟨I⟩)")

# Tabla validacion
tb(s, "Validacion estadistica (NB03)",
   Inches(8.2), Inches(1.15), Inches(4.9), Inches(0.38),
   size=16, bold=True, color=C_ACCENT)

rows_sp = [
    ("Metrica",        "PINN",       "Teorico",       "OK"),
    ("Contraste C",    "1.0253",     "= 1",           ""),
    ("|C - 1|",        "0.025",      "< 0.1",         "cumple"),
    ("KS p-valor",     "< 0.0001",   "> 0.05",        "ver nota"),
    ("Frac. I>2⟨I⟩",   "≈ 0.135",    "e⁻² ≈ 0.1353", "cumple"),
    ("Tiempo total",   "227 s",      "--",            ""),
]
t_top = 1.60
for i, (m, p, t, ok) in enumerate(rows_sp):
    is_hdr = (i == 0)
    is_warn = (i == 3)
    bg_r = (C_ACCENT if is_hdr
            else RGBColor(0xFF, 0xF8, 0xE8) if is_warn
            else (C_LGRAY if i % 2 == 0 else C_WHITE))
    tc   = C_WHITE if is_hdr else C_BODY
    rect(s, Inches(8.2), Inches(t_top+i*0.52), Inches(4.9), Inches(0.5),
         fill=bg_r, line=C_ACCENT2, lw=Pt(0.5))
    for j, (val, cx, cw) in enumerate(zip([m, p, t, ok],
                                           [8.22, 9.6, 10.7, 12.1],
                                           [1.35, 1.08, 1.35, 0.85])):
        col_v = (C_WHITE if is_hdr
                 else C_GREEN if ok == "cumple" and j == 3
                 else C_AMBER if ok == "ver nota" and j == 3
                 else C_BOLD_NUM if i in [1,2] and j == 1
                 else C_BODY)
        tb(s, val, Inches(cx), Inches(t_top+i*0.52+0.09),
           Inches(cw), Inches(0.35), size=12,
           bold=(is_hdr or (i in [1,2] and j==1)), color=col_v,
           align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

# Nota KS
rect(s, Inches(8.2), Inches(4.85), Inches(4.9), Inches(1.45),
     fill=RGBColor(0xFF, 0xF8, 0xE8), line=RGBColor(0xAA, 0x88, 0x00), lw=Pt(1))
tb(s, "Nota sobre KS p < 0.05:",
   Inches(8.35), Inches(4.93), Inches(4.6), Inches(0.35),
   size=13, bold=True, color=C_AMBER)
tb(s, "Con N = 10,000 puntos el test KS tiene alta potencia estadistica: "
      "desviaciones < 1% resultan significativas aunque sean fisicamente irrelevantes. "
      "El criterio principal |C-1| < 0.1 se cumple con margen amplio.",
   Inches(8.35), Inches(5.28), Inches(4.6), Inches(0.9),
   size=11, color=C_BODY)

notes(s, """NB03 VALIDACION - 2 min
La validacion estadistica del speckle sigue el criterio de Goodman (2007).

El criterio mas importante es el contraste C = sigma_I / <I>. Obtenemos C = 1.0253, con |C-1| = 0.025, muy por debajo del umbral de 0.1. Cumple.

La fraccion de pixeles con intensidad mayor a 2 veces la media es 0.135, en excelente concordancia con el valor teorico e^{-2} = 0.1353. Cumple.

El test KS arroja p < 0.0001, lo que formalmente rechaza la hipotesis nula. Sin embargo, esto no indica un fallo fisico: con N = 10,000 puntos, el test tiene altisima potencia estadistica. Desviaciones menores al 1% resultan significativas aunque sean fisicamente irrelevantes.

La distribucion de intensidades sigue muy bien la exponencial teorica, como se ve en la grafica.

En resumen: el modelo genera patrones de speckle estadisticamente validos segun el criterio principal de Goodman.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 14 — COMPARATIVA ESTADO DEL ARTE
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Comparativa con el estado del arte")
slide_num(s, 14)

tb(s, "Referencia: Schoder & Kraxberger (2024) - PINNs convencionales (activacion tanh) para Helmholtz 3D",
   Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.35),
   size=14, italic=True, color=C_ACCENT2)

# Tabla comparativa
hdrs = ["Metodo", "Dimension", "Error L2 (%)", "Factor de mejora"]
rows_c = [
    ("Schoder & Kraxberger (2024)", "1D", "2.490", "--"),
    ("Schoder & Kraxberger (2024)", "2D", "1.910", "--"),
    ("PINN-SIREN (este trabajo)",   "1D", "0.006", "×415"),
    ("PINN-SIREN (este trabajo)",   "2D", "0.171", "×11.2"),
]
cw = [4.5, 1.5, 1.8, 2.4]
cx = [0.4, 4.9, 6.4, 8.2]
t_top = 1.6

for j, (h, w) in enumerate(zip(hdrs, cw)):
    rect(s, Inches(cx[j]), Inches(t_top), Inches(w), Inches(0.52),
         fill=C_ACCENT, line=None)
    tb(s, h, Inches(cx[j]+0.05), Inches(t_top+0.08),
       Inches(w-0.1), Inches(0.35), size=14, bold=True,
       color=C_WHITE, align=PP_ALIGN.CENTER)

for i, (m, d, e, f) in enumerate(rows_c):
    rt = t_top + (i+1)*0.62
    is_ours = (i >= 2)
    bg_r = RGBColor(0xE8, 0xF4, 0xEA) if is_ours else (C_LGRAY if i % 2 == 0 else C_WHITE)
    for j, (val, w) in enumerate(zip([m, d, e, f], cw)):
        rect(s, Inches(cx[j]), Inches(rt), Inches(w), Inches(0.6),
             fill=bg_r, line=C_ACCENT2, lw=Pt(0.5))
        tc = C_GREEN if (is_ours and j in [2, 3]) else C_BODY
        tb(s, val, Inches(cx[j]+0.05), Inches(rt+0.1),
           Inches(w-0.1), Inches(0.42), size=14, bold=is_ours, color=tc,
           align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

# Cajas de factor
rect(s, Inches(0.4), Inches(4.45), Inches(5.5), Inches(1.65),
     fill=C_LGRAY, line=C_ACCENT, lw=Pt(1.5))
tb(s, "×415 en 1D",
   Inches(0.5), Inches(4.55), Inches(5.3), Inches(0.65),
   size=36, bold=True, color=C_BOLD_NUM, align=PP_ALIGN.CENTER)
tb(s, "SIREN vs PINN convencional — Helmholtz 1D",
   Inches(0.5), Inches(5.18), Inches(5.3), Inches(0.75),
   size=14, color=C_BODY, align=PP_ALIGN.CENTER)

rect(s, Inches(6.4), Inches(4.45), Inches(5.5), Inches(1.65),
     fill=C_LGRAY, line=C_ACCENT, lw=Pt(1.5))
tb(s, "×11.2 en 2D",
   Inches(6.5), Inches(4.55), Inches(5.3), Inches(0.65),
   size=36, bold=True, color=C_BOLD_NUM, align=PP_ALIGN.CENTER)
tb(s, "SIREN 2D vs PINN convencional — Helmholtz 2D/3D",
   Inches(6.5), Inches(5.18), Inches(5.3), Inches(0.75),
   size=14, color=C_BODY, align=PP_ALIGN.CENTER)

tb(s, "Nota: la comparacion es indicativa. Schoder & Kraxberger operan en 3D con dominio fisico y activacion tanh; "
      "este trabajo en 1D/2D adimensional con SIREN.",
   Inches(0.4), Inches(6.2), Inches(11.8), Inches(0.35),
   size=11, italic=True, color=C_ACCENT2)

notes(s, """COMPARATIVA - 1.5 min
La comparativa mas relevante es con Schoder y Kraxberger (2024), que resuelven Helmholtz con PINNs de activacion tanh.

Para 1D: su error es 2.49%, el nuestro es 0.006%. Factor de mejora: x415.
Para 2D: su error es 1.91%, el nuestro es 0.171%. Factor de mejora: x11.2.

Es importante ser precisos: la comparacion es indicativa, no directa. Ellos trabajan en 3D con dominio fisico; nosotros en 1D/2D adimensional con SIREN. El factor refleja tanto la ventaja arquitectonica de SIREN como la diferencia en complejidad dimensional.

El mensaje clave: SIREN con optimizacion hibrida Adam + L-BFGS supera significativamente a PINNs convencionales para ecuaciones de onda.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 15 — DISCUSION
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Discusion: Hallazgos principales del avance")
slide_num(s, 15)

# Cajas resumen
boxes_d = [
    ("NB01\nHelmholtz 1D", "L² = 0.006%\nR² = 1.000000\n251 s"),
    ("NB02\nHelmholtz 2D", "L² = 0.171%\nR² > 0.9999\n299 s"),
    ("NB03\nSpeckle optico", "C = 1.025\n|C-1| = 0.025\n227 s"),
]
for i, (tit, met) in enumerate(boxes_d):
    rect(s, Inches(0.4+i*4.2), Inches(1.2), Inches(3.9), Inches(1.6),
         fill=C_ACCENT, line=None)
    tb(s, tit, Inches(0.5+i*4.2), Inches(1.28),
       Inches(3.7), Inches(0.58), size=15, bold=True,
       color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(s, met, Inches(0.5+i*4.2), Inches(1.86),
       Inches(3.7), Inches(0.82), size=15,
       color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

tb(s, "Hallazgos clave:",
   Inches(0.4), Inches(3.0), Inches(12.5), Inches(0.38),
   size=19, bold=True, color=C_ACCENT)

hall = [
    "La misma arquitectura SIREN 5x128 resuelve Helmholtz con onda plana Y con fase rugosa aleatoria sin modificacion -> generalidad del metodo demostrada",
    "Activación sinusoidal SIREN supera ×415 a PINNs convencionales (tanh) para ecuaciones de onda",
    "Optimizacion hibrida Adam + L-BFGS es esencial: Adam evita minimos locales, L-BFGS refina con curvatura",
    "ω₀ = 1.0 y λ_fis = 0.1 son configuraciones no triviales, calibradas empíricamente",
    "Speckle generado cumple el criterio estadístico de Goodman: C ≈ 1, distribución exponencial negativa",
]
bullets(s, hall, Inches(0.4), Inches(3.45), Inches(12.5), Inches(3.8), size=17)

notes(s, """DISCUSION - 1.5 min
Resumiendo los hallazgos principales del primer avance.

Los tres experimentos cumplieron sus objetivos con margen amplio.

El hallazgo mas importante es la generalidad: la misma red SIREN 5x128 resuelve Helmholtz con onda plana y con fase rugosa aleatoria sin ninguna modificacion. Solo cambia la condicion de frontera. Esto sugiere que el metodo es generalizable a cualquier condicion de frontera.

La optimizacion hibrida es esencial: documentamos que L-BFGS puro diverge desde inicializacion aleatoria; Adam puro no alcanza la precision requerida; solo la combinacion converge de forma robusta.

La calibracion de omega_0 y lambda_fis requirio experimentacion y esta documentada como guia para futuras implementaciones.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 16 — TRABAJO FUTURO
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Trabajo futuro y proximos pasos")
slide_num(s, 16)

tb(s, "Experimentos planificados:",
   Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.38),
   size=20, bold=True, color=C_ACCENT)

future = [
    ("NB04", "Benchmark vs FEM (FEniCSx)",
     "Comparar PINN-SIREN vs Metodo de Elementos Finitos. Metrica: S = T_FEM / T_PINN. Cuantificar la aceleracion real."),
    ("NB05", "Medios inhomogeneos (GRIN)",
     "Extender a: ∇²E + k²(r)E = 0 con índice variable. Aplicación: fibras ópticas con gradiente de índice."),
    ("NB06", "Condicion de radiacion de Sommerfeld",
     "Agregar: ∂E/∂n − ikE = 0 en y=1 para eliminar reflexiones espurias y mejorar el test KS del speckle."),
    ("NB07", "Transfer learning entre realizaciones",
     "Reutilizar pesos de una realizacion para la siguiente, reduciendo el costo por nueva fase rugosa."),
]
top = 1.65
for code, title, desc in future:
    rect(s, Inches(0.4), Inches(top), Inches(0.9), Inches(1.1),
         fill=C_ACCENT2, line=None)
    tb(s, code, Inches(0.4), Inches(top+0.07),
       Inches(0.9), Inches(0.45), size=13, bold=True,
       color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, Inches(1.45), Inches(top+0.05),
       Inches(11.3), Inches(0.38), size=17, bold=True, color=C_ACCENT)
    tb(s, desc, Inches(1.45), Inches(top+0.43),
       Inches(11.3), Inches(0.58), size=15, color=C_BODY)
    top += 1.2

notes(s, """TRABAJO FUTURO - 1.5 min
El paso inmediato mas importante es el NB04: el benchmark contra FEniCSx. Esto nos permitira cuantificar el factor de aceleracion S = T_FEM / T_PINN, que es la hipotesis central de la tesis.

NB05 extiende el metodo a medios con indice de refraccion variable: fibras GRIN donde k depende de la posicion.

NB06 incorpora la condicion de radiacion de Sommerfeld en el borde superior, que eliminara las reflexiones espurias que observamos en el test KS.

NB07 explora transfer learning para acelerar la generacion de multiples realizaciones del speckle, que es critico para estudios estadisticos.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 17 — RESUMEN
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Resumen: Que se logro en este primer avance")
slide_num(s, 17)

cols_r = [
    ("Precision demostrada",
     ["Error L2 = 0.006% (1D)", "Error L2 = 0.171% (2D)",
      "Meta < 5% superada por ×29", "R² = 1.000000 en 1D"]),
    ("Speckle estadisticamente valido",
     ["Contraste C = 1.025", "|C - 1| = 0.025 < 0.1",
      "Dist. intensidades: exp. negativa", "Fracción colas: 0.135 ≈ e⁻²"]),
    ("Infraestructura solida",
     ["Codigo modular en src/", "3 notebooks reproducibles",
      "SEED=42 garantiza reproducibilidad", "GPU RTX 5050, CUDA 12.6"]),
]
for i, (tit, items) in enumerate(cols_r):
    x = Inches(0.35 + i*4.3)
    rect(s, x, Inches(1.2), Inches(4.1), Inches(0.55),
         fill=C_ACCENT, line=None)
    tb(s, tit, x+Inches(0.08), Inches(1.25), Inches(3.94), Inches(0.45),
       size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        bg_it = C_LGRAY if j % 2 == 0 else C_WHITE
        rect(s, x, Inches(1.75+j*0.75), Inches(4.1), Inches(0.72),
             fill=bg_it, line=C_ACCENT2, lw=Pt(0.5))
        tb(s, item, x+Inches(0.08), Inches(1.8+j*0.75),
           Inches(3.94), Inches(0.6), size=15, color=C_BODY)

rect(s, Inches(0.35), Inches(4.85), Inches(12.6), Inches(0.7),
     fill=C_LGRAY, line=C_ACCENT, lw=Pt(1.5))
tb(s, "Pregunta de investigacion respondida parcialmente: la PINN-SIREN simula "
      "speckle con C ≈ 1  ✓   |   Falta: benchmark FEM para cuantificar aceleración",
   Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.5),
   size=15, italic=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

notes(s, """RESUMEN - 1 min
En sintesis, este primer avance logro tres cosas:

Primero: demostrar la precision del metodo, con errores muy por debajo del umbral.

Segundo: validar estadisticamente el speckle generado segun el criterio de Goodman.

Tercero: construir la infraestructura computacional reproducible.

La pregunta central de la tesis tiene respuesta positiva en cuanto a calidad estadistica. Lo que falta es el benchmark FEM para cuantificar la aceleracion respecto al metodo clasico.""")

# ──────────────────────────────────────────────────────────────────────────
# SLIDE 18 — CIERRE
# ──────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank())
bg(s, C_ACCENT)
rect(s, 0, 0, SW, Inches(0.12), fill=C_WHITE)
rect(s, 0, SH-Inches(0.12), SW, Inches(0.12), fill=C_WHITE)

tb(s, "Gracias por su atencion",
   Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.85),
   size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(1.8), Inches(2.0), Inches(9.7), Inches(0.03),
     fill=RGBColor(0xCC, 0xDD, 0xEE))

tb(s, "Roberto Hernandez Estrada",
   Inches(0.5), Inches(2.15), Inches(12.3), Inches(0.48),
   size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "robertohernandezestrd@gmail.com   |   Matricula: 252H21004",
   Inches(0.5), Inches(2.62), Inches(12.3), Inches(0.35),
   size=14, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(1.8), Inches(3.1), Inches(9.7), Inches(0.03),
     fill=RGBColor(0x80, 0xA0, 0xBB))

tb(s, "Agradecimientos",
   Inches(0.5), Inches(3.28), Inches(12.3), Inches(0.38),
   size=18, bold=True, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

acks = [
    "Dr. Jose Adan Hernandez Nolasco  -  Director de tesis",
    "Dr. Oscar Alberto Chavez Bosquez  -  Codirector de tesis",
    "DACYTI-UJAT  -  Programa de Maestria en Ciencias de la Computacion",
]
for k, ack in enumerate(acks):
    tb(s, ack, Inches(1.5), Inches(3.75+k*0.5), Inches(10.3), Inches(0.42),
       size=17, color=C_WHITE, align=PP_ALIGN.CENTER)

tb(s, "Python 3.14  |  PyTorch 2.4  |  CUDA 12.6  |  GPU NVIDIA RTX 5050",
   Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.38),
   size=13, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(1.8), Inches(6.1), Inches(9.7), Inches(0.03),
     fill=RGBColor(0x80, 0xA0, 0xBB))

tb(s, "Preguntas?",
   Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
   size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

notes(s, """CIERRE - 0.5 min
Muchas gracias por su atencion. Quedo a sus ordenes para preguntas.

Posibles preguntas:
- Por que SIREN y no transformers? SIREN tiene derivadas analiticas exactas necesarias para nabla^2
- Por que L-BFGS converge en pocas iteraciones? La buena inicializacion de Sitzmann parte de un punto favorable
- Que tan escalable es a 3D? NB05 planificado, principal reto computacional
- El speckle generado es util en la practica? Si, como sustituto computacional rapido para estudios estadisticos""")

# ── Guardar ────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)
print(f"Presentacion guardada en: {OUT_PATH}")
print(f"Total de slides: {len(prs.slides)}")
for i, fig in enumerate([FIG_1D_RES, FIG_1D_MET, FIG_2D_RES, FIG_2D_MET, FIG_SP_RES, FIG_SP_EST]):
    estado = "OK" if os.path.exists(fig) else "NO ENCONTRADA"
    print(f"  Figura {i+1}: {os.path.basename(fig)} -> {estado}")
