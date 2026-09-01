"""
add_code_slides.py
Agrega 2 slides de implementación computacional a slide_new.pptx
y las inserta después de la diapositiva METODOLOGÍA (índice 12).

Slide A: Pipeline de implementación (flujo real del código)
Slide B: Arquitectura del código (src/ módulos)
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_IN   = r'D:\Tesis_Maestria\slides\presentacion_mcc_roberto_hernandez_estrada.pptx'
PPTX_PATH = r'D:\Tesis_Maestria\slides\slide_new.pptx'

import shutil
shutil.copy2(PPTX_IN, PPTX_PATH)
prs = Presentation(PPTX_PATH)
print(f"Cargado desde: {PPTX_IN}")
print(f"Total slides base: {len(prs.slides)}")

# ── Colores ─────────────────────────────────────────────
C_DARK  = RGBColor(0x1A, 0x1A, 0x2E)   # fondo oscuro código
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE  = RGBColor(0x1A, 0x52, 0x7A)
C_GREEN = RGBColor(0x1A, 0x6B, 0x3A)
C_RED   = RGBColor(0xC0, 0x39, 0x4B)
C_PURP  = RGBColor(0x6B, 0x2B, 0x9A)
C_ORNG  = RGBColor(0xB8, 0x47, 0x0F)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_LITE  = RGBColor(0x88, 0x88, 0x88)
C_TXT   = RGBColor(0x22, 0x22, 0x22)
C_MONO  = RGBColor(0x0D, 0x47, 0xA1)   # azul código
C_CODE_BG = RGBColor(0xF0, 0xF4, 0xF8) # fondo caja código

def hex6(rgb):
    r, g, b = rgb
    return f'{r:02X}{g:02X}{b:02X}'

def _make_para(text, bold=False, size_pt=None, color=None,
               align=None, italic=False, mono=False):
    p = etree.Element(qn('a:p'))
    if align:
        pPr = etree.SubElement(p, qn('a:pPr'))
        m = {PP_ALIGN.CENTER: 'ctr', PP_ALIGN.LEFT: 'l', PP_ALIGN.RIGHT: 'r'}
        if align in m:
            pPr.set('algn', m[align])
    if not text:
        ep = {'lang': 'es-MX'}
        if size_pt:
            ep['sz'] = str(int(size_pt * 100))
        etree.SubElement(p, qn('a:endParaRPr'), attrib=ep)
        return p
    r = etree.SubElement(p, qn('a:r'))
    rpr = {'lang': 'es-MX', 'dirty': '0'}
    if bold:   rpr['b'] = '1'
    if italic: rpr['i'] = '1'
    if size_pt: rpr['sz'] = str(int(size_pt * 100))
    rPr = etree.SubElement(r, qn('a:rPr'), attrib=rpr)
    if mono:
        latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', 'Consolas')
    if color:
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': hex6(color)})
    t = etree.SubElement(r, qn('a:t'))
    t.text = text
    return p

def _set_txbody(tf, lines):
    txb = tf._txBody
    for p in txb.findall(qn('a:p')):
        txb.remove(p)
    for line in lines:
        text = line[0]
        bold = line[1] if len(line) > 1 else False
        size = line[2] if len(line) > 2 else None
        color = line[3] if len(line) > 3 else None
        align = line[4] if len(line) > 4 else None
        italic = line[5] if len(line) > 5 else False
        mono = line[6] if len(line) > 6 else False
        txb.append(_make_para(text, bold, size, color, align, italic, mono))

def add_box(slide, l, t, w, h, lines, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = wrap
    _set_txbody(box.text_frame, lines)
    return box

def add_rect(slide, l, t, w, h, fill_color, lines=None, wrap=True):
    """Agrega un rectángulo con relleno de color y texto opcional."""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # sin borde
    if lines:
        shape.text_frame.word_wrap = wrap
        _set_txbody(shape.text_frame, lines)
    return shape

def set_notes(slide, text):
    if not slide.has_notes_slide:
        _ = slide.notes_slide
    for shape in slide.notes_slide.shapes:
        if shape.has_text_frame:
            try:
                shape.text_frame.text = text
                return
            except Exception:
                continue

def insert_at(prs, position):
    """Mueve el último slide (recién añadido) a la posición indicada (0-based)."""
    xml_slides = prs.slides._sldIdLst
    last = xml_slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(position, last)

# ── Layout en blanco ────────────────────────────────────
blank = None
for lay in prs.slide_layouts:
    if lay.name.lower() in ('blank', 'en blanco', 'vacío'):
        blank = lay
        break
if blank is None:
    blank = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════
# SLIDE A — PIPELINE DE IMPLEMENTACIÓN
# ════════════════════════════════════════════════════════
sA = prs.slides.add_slide(blank)

# Título
add_box(sA, 0.3, 0.12, 12.7, 0.6, [
    ("IMPLEMENTACIÓN: Pipeline de Entrenamiento",
     True, 22, C_BLUE, PP_ALIGN.CENTER)
])
add_box(sA, 0.3, 0.75, 12.7, 0.35, [
    ("Flujo real del código — src/losses.py · src/training.py · torch.autograd",
     False, 13, C_GRAY, PP_ALIGN.CENTER)
])

# ── Caja 1: Entrada / Muestreo ──────────────────────────
add_rect(sA, 0.3, 1.25, 3.8, 1.5,
    RGBColor(0xE3, 0xF2, 0xFD),
    [
        ("① Entrada y Muestreo", True, 13, C_BLUE, None),
        ("", False, 5, None, None),
        ("LHS(n=3000) → xy ∈ [0,1]²",     False, 11, C_TXT, None, False, True),
        ("xy_bc → condiciones de frontera", False, 11, C_TXT, None, False, True),
        ("SEED=42 → reproducibilidad",      False, 11, C_TXT, None, False, True),
    ]
)

# Flecha →
add_box(sA, 4.15, 1.8, 0.5, 0.5, [("→", True, 20, C_GRAY, PP_ALIGN.CENTER)])

# ── Caja 2: Forward Pass ──────────────────────────────
add_rect(sA, 4.7, 1.25, 3.8, 1.5,
    RGBColor(0xE8, 0xF5, 0xE9),
    [
        ("② Forward Pass — models.py", True, 13, C_GREEN, None),
        ("", False, 5, None, None),
        ("PINN_2D_SIREN.forward(xy)",     False, 11, C_TXT, None, False, True),
        ("  5 × Linear(n,128) + Sine(ω₀)", False, 11, C_TXT, None, False, True),
        ("→ (E_real, E_imag) ∈ ℝ²",      False, 11, C_TXT, None, False, True),
    ]
)

# Flecha →
add_box(sA, 8.55, 1.8, 0.5, 0.5, [("→", True, 20, C_GRAY, PP_ALIGN.CENTER)])

# ── Caja 3: Residuo Helmholtz ─────────────────────────
add_rect(sA, 9.1, 1.25, 3.8, 1.5,
    RGBColor(0xF3, 0xE5, 0xF5),
    [
        ("③ Residuo Helmholtz — losses.py", True, 13, C_PURP, None),
        ("", False, 5, None, None),
        ("autograd.grad(E, xy) → ∇E",        False, 11, C_TXT, None, False, True),
        ("autograd.grad(∇E, xy) → ∇²E",      False, 11, C_TXT, None, False, True),
        ("R = ∇²E + k²·E  [sin malla]",      False, 11, C_TXT, None, False, True),
    ]
)

# Flechas verticales hacia abajo ↓
add_box(sA, 5.3, 2.8, 1.0, 0.5, [("↓", True, 20, C_GRAY, PP_ALIGN.CENTER)])
add_box(sA, 10.7, 2.8, 1.0, 0.5, [("↓", True, 20, C_GRAY, PP_ALIGN.CENTER)])

# ── Caja 4: Función de pérdida ────────────────────────
add_rect(sA, 3.5, 3.3, 6.3, 1.3,
    RGBColor(0xFF, 0xF3, 0xE0),
    [
        ("④ Función de pérdida — losses.py::pinn_loss_2d()", True, 13, C_ORNG, None),
        ("", False, 5, None, None),
        ("ℒ = MSE(E_pred_bc − E_bc)  +  λ·[MSE(R_real²) + MSE(R_imag²)]",
         False, 12, C_TXT, None, False, True),
        ("    ←── ℒ_frontera ───→          ←────────── λ = 0.1 ──────────────→",
         False, 10, C_GRAY, None, True),
    ]
)

# Flecha ↓
add_box(sA, 6.3, 4.65, 1.0, 0.5, [("↓", True, 20, C_GRAY, PP_ALIGN.CENTER)])

# ── Caja 5: Optimización ─────────────────────────────
add_rect(sA, 1.8, 5.2, 9.7, 1.3,
    RGBColor(0xFC, 0xE4, 0xEC),
    [
        ("⑤ Optimización — training.py::train_adam_lbfgs()", True, 13, C_RED, None),
        ("", False, 5, None, None),
        ("Adam(lr=1e-3) → early stopping (patience=800, δ<1e-7) → best_state guardado",
         False, 11, C_TXT, None, False, True),
        ("L-BFGS(strong_wolfe, history=100, max_iter=1000) → convergencia cuadrática",
         False, 11, C_TXT, None, False, True),
    ]
)

# Flecha ↓
add_box(sA, 6.3, 6.55, 1.0, 0.5, [("↓", True, 20, C_GRAY, PP_ALIGN.CENTER)])

# ── Caja 6: Métricas y guardado ───────────────────────
add_rect(sA, 3.0, 7.05, 7.3, 0.7,
    RGBColor(0xE8, 0xF5, 0xE9),
    [
        ("⑥  utils.py::l2_rel() → L²  |  R²  |  Pearson  →  save_model() → .pt",
         False, 12, C_GREEN, None, False, True),
    ]
)

set_notes(sA,
    "PIPELINE DE IMPLEMENTACIÓN - 1:30 min. "
    "La clave computacional es que el laplaciano se calcula con torch.autograd.grad() "
    "dos veces encadenadas — sin diferencias finitas, sin malla. "
    "El resultado es exacto a la precisión de float32. "
    "training.py implementa early stopping con guardado del mejor estado "
    "y L-BFGS con closure pattern para convergencia cuadrática."
)

insert_at(prs, 13)   # insertar después de METODOLOGÍA (índice 12)
print("[A] Slide 'Pipeline de Implementación' insertado en posición 14")


# ════════════════════════════════════════════════════════
# SLIDE B — ARQUITECTURA DEL CÓDIGO
# ════════════════════════════════════════════════════════
sB = prs.slides.add_slide(blank)

# Título
add_box(sB, 0.3, 0.12, 12.7, 0.6, [
    ("ARQUITECTURA DEL CÓDIGO — Módulos src/",
     True, 22, C_BLUE, PP_ALIGN.CENTER)
])
add_box(sB, 0.3, 0.75, 12.7, 0.35, [
    ("Diseño modular: 4 módulos especializados + 2 notebooks reproducibles",
     False, 13, C_GRAY, PP_ALIGN.CENTER)
])

# ── Columna izquierda: módulos src/ ───────────────────
mod_data = [
    ("models.py", C_GREEN,
     ["Sine(ω₀)  |  PINN_1D_SIREN  |  PINN_2D_SIREN",
      "_init_weights():  W~U(−√6/n/ω₀, +√6/n/ω₀)",
      "forward(xy) → (E_real, E_imag)",
      "save() / load()  para reutilizar entre NB"]),
    ("losses.py", C_PURP,
     ["helmholtz_residual_1d / _2d",
      "∇²E = autograd.grad(∇E, xy)  [sin malla]",
      "pinn_loss_1d / _2d",
      "ℒ = ℒ_BC  +  λ·ℒ_física"]),
    ("training.py", C_RED,
     ["train_adam_lbfgs(model, loss_fn, config)",
      "Adam  →  early stop (patience, best_state)",
      "L-BFGS  →  closure + strong_wolfe",
      "Retorna history, iter_lbfgs, tiempos"]),
    ("utils.py", C_ORNG,
     ["l2_rel(pred, exact)  ← métrica principal",
      "get_figures_dir() / get_models_dir()",
      "save_model(model, 'nb02_helmholtz2d')",
      "load_model(model, 'nb02_...')  → NB03"]),
]

ypos = 1.3
for modname, color, items in mod_data:
    # Encabezado del módulo
    add_rect(sB, 0.3, ypos, 5.6, 0.38,
        color,
        [("src/" + modname, True, 13, C_WHITE, None, False, True)]
    )
    # Contenido
    lines = []
    for item in items:
        lines.append((item, False, 11, C_TXT, None, False, True))
    add_rect(sB, 0.3, ypos + 0.38, 5.6, 0.72,
        RGBColor(0xF8, 0xF9, 0xFA),
        lines
    )
    ypos += 1.2

# ── Columna derecha: notebooks + flujo ───────────────
add_box(sB, 6.3, 1.3, 6.7, 0.4, [
    ("Notebooks — flujo de uso", True, 14, C_BLUE, None)
])

nb_data = [
    ("NB01", "Helmholtz 1D",    "models + losses (loop inline)", "SEED=42", C_GREEN,
     "results/models/nb01_helmholtz1d.pt"),
    ("NB02", "Helmholtz 2D",    "models + losses (loop inline)", "SEED=42", C_BLUE,
     "results/models/nb02_helmholtz2d.pt"),
    ("NB03", "Speckle óptico",  "models + training (reutiliza)", "Pendiente", C_ORNG,
     "results/models/nb03_speckle.pt"),
]

ypos2 = 1.8
for nb, name, modules, seed, color, outfile in nb_data:
    add_rect(sB, 6.3, ypos2, 6.7, 1.05,
        RGBColor(0xF0, 0xF4, 0xF8),
        [
            (f"{nb} — {name}", True, 13, color, None),
            ("", False, 4, None, None),
            (f"from src  →  {modules}", False, 11, C_TXT, None, False, True),
            (f"→  {outfile}",            False, 10, C_LITE, None, False, True),
        ]
    )
    ypos2 += 1.15

# Stack tecnológico
add_rect(sB, 6.3, 5.25, 6.7, 0.95,
    RGBColor(0x1A, 0x1A, 0x2E),
    [
        ("Stack tecnológico", True, 12, C_WHITE, None),
        ("", False, 4, None, None),
        ("Python 3.11  |  PyTorch 2.4  |  CUDA 12.6", False, 11,
         RGBColor(0xA0, 0xC4, 0xFF), None, False, True),
        ("GPU: NVIDIA RTX 5050  |  SEED=42 global", False, 11,
         RGBColor(0xA0, 0xC4, 0xFF), None, False, True),
    ]
)

# Nota de diseño
add_box(sB, 0.3, 6.35, 12.7, 0.55, [
    ("Decisión de diseño: training loop inline en NB01/NB02 (pedagógico) "
     "→ modularizado en training.py para NB03+",
     False, 11, C_GRAY, PP_ALIGN.CENTER, True)
])

set_notes(sB,
    "ARQUITECTURA DEL CÓDIGO - 1:00 min. "
    "Cuatro módulos especializados: models (arquitectura SIREN + inicialización), "
    "losses (residuo Helmholtz con autograd, función de pérdida), "
    "training (Adam early-stop + L-BFGS closure), "
    "utils (métricas, save/load para reusar pesos entre notebooks). "
    "NB01 y NB02 tienen el loop inline por legibilidad pedagógica; "
    "training.py factoriza ese loop para NB03+."
)

insert_at(prs, 14)   # insertar después del pipeline (posición 14)
print("[B] Slide 'Arquitectura del Código' insertado en posición 15")

# ── Guardar ──────────────────────────────────────────
prs.save(PPTX_PATH)
print(f"\nTotal slides: {len(prs.slides)}")
print(f"✅  Guardado en: {PPTX_PATH}")
