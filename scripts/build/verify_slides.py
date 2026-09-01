import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
prs = Presentation(r'D:\Tesis_Maestria\slides\slide_new.pptx')
print(f'Total slides: {len(prs.slides)}')
print()
for i in range(10, 18):
    s = prs.slides[i]
    titles = [sh.text_frame.text[:70] for sh in s.shapes
              if sh.has_text_frame and sh.text_frame.text.strip()]
    first = titles[0] if titles else '--vacia--'
    print(f'  Slide {i+1}: {first}')
