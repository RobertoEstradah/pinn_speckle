import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

prs = Presentation(r'D:\Tesis_Maestria\slides\MANUAL\slide_maestria_roberto_hernandez_estrada_MANUAL.pptx')

print(f'Total slides: {len(prs.slides)}')
print(f'Slide dimensions: {prs.slide_width.inches:.1f}" x {prs.slide_height.inches:.1f}"')
print()

for i, slide in enumerate(prs.slides, 1):
    print(f'## Slide {i}')
    for shape in slide.shapes:
        if shape.has_text_frame:
            name = shape.name
            texts = []
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)
            if texts:
                print(f'  [{name}]')
                for t in texts:
                    print(f'    {t}')
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f'  [NOTES]: {notes[:300]}')
    print()
