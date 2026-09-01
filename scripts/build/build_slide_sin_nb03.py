#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Genera presentacion PPTX para el Primer Avance de Tesis - VERSION SIN NB03
Simulacion Acelerada de Speckle Optico mediante PINNs
Roberto Hernandez Estrada - UJAT DACYTI
Slides: 16 (se eliminaron slides 12 y 13 de NB03)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Rutas ──────────────────────────────────────────────────────────────────
BASE_DIR = r"D:\Tesis_Maestria"
FIG_DIR  = os.path.join(BASE_DIR, "results", "figures")
OUT_PATH = os.path.join(BASE_DIR, "slides", "slide_maestria_roberto_hernandez_estrada.pptx")

FIG_1D_RES = os.path.join(FIG_DIR, "resultados_pinn_1d.png")
FIG_1D_MET = os.path.join(FIG_DIR, "metricas_adicionales_1d.png")
FIG_2D_RES = os.path.join(FIG_DIR, "resultados_pinn_2d.png")
FIG_2D_MET = os.path.join(FIG_DIR, "metricas_adicionales_2d.png")

# ── Paleta ─────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0xFA, 0xFA, 0xF8)
C_TITLE    = RGBColor(0x1A, 0x1A, 0x2E)
C_BODY     = RGBColor(0x2C, 0x2C, 0x2C)
C_ACCENT   = RGBColor(0x34, 0x5C, 0x8C)
C_ACCENT2  = RGBColor(0x7A, 0x92, 0xAA)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY    = RGBColor(0xF0, 0xF2, 0xF5)
C_MGRAY    = RGBColor(0xCC, 0xCC, 0xCC)
C_BOLD_NUM = RGBColor(0x1B, 0x4F, 0x72)
C_GREEN    = RGBColor(0x1A, 0x5C, 0x2A)
C_AMBER    = RGBColor(0x99, 0x55, 0x00)
C_SLIDE_N  = RGBColor(0xAA, 0xAA, 0xAA)

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

def blank():
    return prs.slide_layouts[6]

def bg(slide, color=None):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color or C_BG

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, size=18, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name  = "Calibri"
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color or C_BODY
    return box

def header(slide, title, sub=None, bar=None):
    bc = bar or C_ACCENT
    rect(slide, 0, 0, SW, Inches(1.05), fill=bc)
    tb(slide, title, Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.78),
       size=27, bold=True, color=C_WHITE)
    if sub:
        tb(slide, sub, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.25),
           size=12, italic=True, color=RGBColor(0xBB, 0xCF, 0xE5))

def slide_num(slide, n):
    tb(slide, str(n), Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
       size=11, color=C_SLIDE_N, align=PP_ALIGN.RIGHT)

def bullets(slide, items, l, t, w, h, size=18, color=None, bchar="  "):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        p.space_after  = Pt(2)
        r = p.add_run()
        if isinstance(item, tuple):
            txt, sz, bold, col = item
            r.text = txt
            r.font.size  = Pt(sz)
            r.font.bold  = bold
            r.font.color.rgb = col or color or C_BODY
        else:
            r.text = bchar + item
            r.font.size  = Pt(size)
            r.font.color.rgb = color or C_BODY
        r.font.name = "Calibri"

def metric_box(slide, label, value, unit, l, t, w=Inches(2.2), h=Inches(1.4)):
    rect(slide, l, t, w, h, fill=C_LGRAY, line=C_ACCENT, lw=Pt(1.5))
    tb(slide, value,  l+Inches(0.08), t+Inches(0.08),
       w-Inches(0.16), Inches(0.62), size=30, bold=True,
       color=C_BOLD_NUM, align=PP_ALIGN.CENTER)
    tb(slide, unit,   l+Inches(0.08), t+Inches(0.7),
       w-Inches(0.16), Inches(0.28), size=12,
       color=C_ACCENT2, align=PP_ALIGN.CENTER)
    tb(slide, label,  l+Inches(0.08), t+Inches(0.98),
       w-Inches(0.16), Inches(0.32), size=12, bold=True,
       color=C_BODY, align=PP_ALIGN.CENTER)

def figure(slide, path, l, t, w, h, cap=None):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, width=w, height=h)
    else:
        rect(slide, l, t, w, h, fill=C_LGRAY, line=C_MGRAY, lw=Pt(1))
        tb(slide, "[" + os.path.basename(path) + "]",
           l+Inches(0.1), t+h/2-Inches(0.2), w-Inches(0.2), Inches(0.4),
           size=11, color=C_MGRAY, align=PP_ALIGN.CENTER)
    if cap:
        tb(slide, cap, l, t+h+Inches(0.02), w, Inches(0.28),
           size=10, italic=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s, C_ACCENT)
rect(s, 0, 0, SW, Inches(0.12), fill=C_WHITE)
rect(s, 0, SH-Inches(0.12), SW, Inches(0.12), fill=C_WHITE)

tb(s, "UNIVERSIDAD JUAREZ AUTONOMA DE TABASCO",
   Inches(0.8), Inches(0.22), Inches(11.7), Inches(0.45),
   size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Division Academica de Ciencias y Tecnologias de la Informacion (DACYTI)",
   Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.38),
   size=12, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(1.5), Inches(1.18), Inches(10.3), Inches(0.025),
     fill=RGBColor(0xBB, 0xCC, 0xDD))

tb(s, "Simulacion Acelerada de Speckle Optico",
   Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.75),
   size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "mediante Redes Neuronales Informadas por Fisica con Activacion Sinusoidal",
   Inches(0.5), Inches(2.08), Inches(12.3), Inches(0.6),
   size=20, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(2.0), Inches(2.82), Inches(9.3), Inches(0.03),
     fill=RGBColor(0xCC, 0xDD, 0xEE))

tb(s, "PRIMER AVANCE  -  COLOQUIO DE MAESTRIA",
   Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.38),
   size=14, bold=True, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

tb(s, "Roberto Hernandez Estrada",
   Inches(0.5), Inches(3.48), Inches(12.3), Inches(0.48),
   size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Maestria en Ciencias de la Computacion  |  Matricula: 252H21004",
   Inches(0.5), Inches(3.94), Inches(12.3), Inches(0.35),
   size=13, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(2.0), Inches(4.40), Inches(9.3), Inches(0.025),
     fill=RGBColor(0x80, 0xA0, 0xBB))

tb(s, "Director: Dr. Jose Adan Hernandez Nolasco",
   Inches(0.5), Inches(4.52), Inches(12.3), Inches(0.35),
   size=15, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Codirector: Dr. Oscar Alberto Chavez Bosquez",
   Inches(0.5), Inches(4.86), Inches(12.3), Inches(0.35),
   size=15, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "Cunduacan, Tabasco  -  Junio 2026",
   Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.38),
   size=13, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

notes(s, "PORTADA - ~0:45 min\nBuenos dias. Mi nombre es Roberto Hernandez Estrada, alumno del programa de Maestria en Ciencias de la Computacion de la DACYTI-UJAT.\nPresento el primer avance de mi tesis bajo la direccion del Dr. Hernandez Nolasco y codireccion del Dr. Chavez Bosquez.\nEn los proximos 20 minutos: motivacion, metodologia y resultados de NB01 y NB02.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA (sin NB03)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Contenido de la presentacion")
slide_num(s, 2)

agenda = [
    "1.  Motivacion y el problema computacional",
    "2.  Objetivos del avance",
    "3.  Marco teorico: Speckle optico y ecuacion de Helmholtz",
    "4.  PINNs y arquitectura SIREN",
    "5.  Metodologia: muestreo, perdida, optimizacion",
    "6.  Resultados NB01 - Validacion Helmholtz 1D",
    "7.  Resultados NB02 - Helmholtz 2D campo complejo",
    "8.  Comparativa con el estado del arte",
    "9.  Trabajo futuro y conclusiones preliminares",
]
box = s.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(6.0))
tf  = box.text_frame
tf.word_wrap = True
for i, item in enumerate(agenda):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(10)
    r = p.add_run()
    r.text = item
    r.font.name = "Calibri"
    r.font.size = Pt(21)
    r.font.color.rgb = C_BODY

notes(s, "AGENDA - 0:30 min\nEstructura en tres bloques: contexto y teoria, dos experimentos con resultados, comparativa y proximos pasos.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MOTIVACION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Motivacion: Por que simular speckle optico?")
slide_num(s, 3)

tb(s, "El speckle optico aparece en:",
   Inches(0.4), Inches(1.15), Inches(5.8), Inches(0.38),
   size=18, bold=True, color=C_ACCENT)
apps = [
    "Metrologia optica y deteccion de defectos superficiales",
    "Tomografia de coherencia optica (OCT)",
    "Imagenologia biomedica",
    "Criptografia optica",
    "Comunicaciones por fibra optica",
]
bullets(s, apps, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.2), size=18)

rect(s, Inches(6.8), Inches(1.1), Inches(6.1), Inches(2.7),
     fill=C_LGRAY, line=C_ACCENT, lw=Pt(1))
tb(s, "El reto computacional con metodos clasicos (FEM):",
   Inches(7.0), Inches(1.2), Inches(5.7), Inches(0.38),
   size=17, bold=True, color=C_ACCENT)
probs = [
    "Malla debe tener resolucion < longitud de onda",
    "Millones de incognitas por realizacion",
    "Cada nuevo perfil de rugosidad: resolver desde cero",
    "Para estudios estadisticos (miles de realizaciones): costo intratable",
]
bullets(s, probs, Inches(7.0), Inches(1.65), Inches(5.7), Inches(2.0), size=16)

rect(s, Inches(6.8), Inches(3.95), Inches(6.1), Inches(2.0),
     fill=C_ACCENT, line=None)
tb(s, "Propuesta: PINN libre de malla",
   Inches(7.0), Inches(4.05), Inches(5.7), Inches(0.38),
   size=17, bold=True, color=C_WHITE)
tb(s, "Incorporar la fisica directamente en la funcion de perdida "
      "de una red neuronal. Una vez entrenada, evalua el campo "
      "completo en ~1.3 ms (GPU RTX 5050).",
   Inches(7.0), Inches(4.45), Inches(5.7), Inches(1.35),
   size=16, color=C_WHITE)

notes(s, "MOTIVACION - 1:30 min\nEl speckle optico aparece en aplicaciones criticas. El FEM es preciso pero muy costoso: malla sub-longitud-de-onda, millones de incognitas, y cada perfil de rugosidad exige resolver desde cero. Para estudios estadisticos esto es intratable.\nPropuesta: PINN libre de malla, evaluable en milisegundos una vez entrenada.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OBJETIVOS (NB03 como pendiente)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Objetivos de investigacion")
slide_num(s, 4)

tb(s, "Objetivo general",
   Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.38),
   size=20, bold=True, color=C_ACCENT)
rect(s, Inches(0.4), Inches(1.58), Inches(12.5), Inches(0.9),
     fill=C_LGRAY, line=C_ACCENT2, lw=Pt(1))
tb(s, "Desarrollar y validar una PINN-SIREN que simule speckle optico resolviendo la "
      "ecuacion de Helmholtz 2D con error L2 < 5% y contraste C ≈ 1.",
   Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.75),
   size=18, color=C_BODY)

rect(s, Inches(0.4), Inches(2.6), Inches(12.5), Inches(0.025), fill=C_MGRAY)

tb(s, "Estado de los experimentos en este avance",
   Inches(0.4), Inches(2.7), Inches(12.5), Inches(0.38),
   size=20, bold=True, color=C_ACCENT)

objs = [
    ("NB01", "Validar la arquitectura SIREN con solucion analitica exacta (Helmholtz 1D)", C_GREEN,   "Completado"),
    ("NB02", "Extender al caso 2D con campo electrico complejo y muestreo LHS",            C_GREEN,   "Completado"),
    ("NB03", "Generar patrones de speckle con validacion estadistica (criterio de Goodman)",C_AMBER,  "Pendiente"),
    ("Meta", "Error L2 < 5%  |  Contraste |C - 1| < 0.1",                                 C_BOLD_NUM, ""),
]
top = 3.18
for code, desc, col, estado in objs:
    rect(s, Inches(0.4), Inches(top), Inches(0.95), Inches(0.57), fill=col, line=None)
    tb(s, code, Inches(0.4), Inches(top+0.03),
       Inches(0.95), Inches(0.5), size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(s, desc, Inches(1.5), Inches(top+0.08),
       Inches(9.8), Inches(0.42), size=18, color=C_BODY)
    if estado:
        tc_e = C_GREEN if estado == "Completado" else C_AMBER
        tb(s, estado, Inches(11.4), Inches(top+0.08),
           Inches(1.5), Inches(0.42), size=14, bold=True, color=tc_e, align=PP_ALIGN.CENTER)
    top += 0.75

notes(s, "OBJETIVOS - 0:45 min\nEl objetivo general es la PINN-SIREN para speckle. En este avance: NB01 y NB02 completados con margen amplio. NB03 es el siguiente paso inmediato.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HELMHOLTZ Y SPECKLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Marco teorico: Ecuacion de Helmholtz y speckle optico")
slide_num(s, 5)

tb(s, "Ecuacion de Helmholtz escalar 2D:",
   Inches(0.4), Inches(1.15), Inches(6.0), Inches(0.38), size=18, bold=True, color=C_ACCENT)
rect(s, Inches(0.4), Inches(1.58), Inches(6.0), Inches(1.6),
     fill=C_LGRAY, line=C_ACCENT2, lw=Pt(1))
tb(s, "∇²E(r) + k²E(r) = 0",
   Inches(0.6), Inches(1.72), Inches(5.6), Inches(0.5),
   size=24, bold=True, color=C_BOLD_NUM, align=PP_ALIGN.CENTER)
tb(s, "k = 2π/λ,   λ = 638 nm (laser diodo rojo)\n"
      "Dominio: Ω = [0,1]²  (adimensional)\n"
      "Campo complejo: E = E_real + i·E_imag",
   Inches(0.6), Inches(2.2), Inches(5.6), Inches(0.85), size=14, color=C_BODY)

tb(s, "La red predice ambas componentes independientemente:",
   Inches(0.4), Inches(3.3), Inches(6.0), Inches(0.35), size=16, color=C_BODY)
rect(s, Inches(0.4), Inches(3.7), Inches(6.0), Inches(0.85),
     fill=C_LGRAY, line=C_ACCENT2, lw=Pt(1))
tb(s, "∇²E_real + k²E_real = 0\n"
      "∇²E_imag + k²E_imag = 0",
   Inches(0.6), Inches(3.78), Inches(5.6), Inches(0.68),
   size=17, bold=True, color=C_BOLD_NUM)

tb(s, "Estadistica del speckle completamente desarrollado:",
   Inches(7.0), Inches(1.15), Inches(6.0), Inches(0.38), size=18, bold=True, color=C_ACCENT)
hechos = [
    "Intensidad: I = |E|² = E²_real + E²_imag",
    "Distribucion teorica: exponencial negativa",
    "p(I) = (1/⟨I⟩)·exp(−I/⟨I⟩)",
    "Contraste: C = σ_I / ⟨I⟩ = 1  (speckle puro)",
    "Criterio validacion: |C − 1| < 0.1  (Goodman 2007)",
]
bullets(s, hechos, Inches(7.0), Inches(1.6), Inches(6.0), Inches(2.8), size=17)

rect(s, Inches(7.0), Inches(4.55), Inches(6.0), Inches(2.2),
     fill=C_ACCENT2, line=None)
tb(s, "Condicion de frontera para speckle (NB03 - proximo paso):",
   Inches(7.2), Inches(4.65), Inches(5.6), Inches(0.38),
   size=14, bold=True, color=C_WHITE)
tb(s, "E(x,0) = e^(iφ(x)),   φ(x) ∼ U(0, 2π)\n"
      "N_φ = 256 puntos en y=0\n"
      "Bordes x=0, x=1, y=1: libres (sin Dirichlet)\n"
      "→ Modela superficie rugosa iluminada por laser",
   Inches(7.2), Inches(5.05), Inches(5.6), Inches(1.55), size=15, color=C_WHITE)

notes(s, "MARCO TEORICO - 1:30 min\nLa fisica esta gobernada por la ecuacion de Helmholtz. Descomponemos el campo complejo en partes real e imaginaria.\nPara el speckle: el criterio de Goodman es el contraste C = 1. La condicion de frontera de fase aleatoria se implementara en NB03.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PINNs Y SIREN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Physics-Informed Neural Networks y arquitectura SIREN")
slide_num(s, 6)

tb(s, "Funcion de perdida PINN:",
   Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.38), size=19, bold=True, color=C_ACCENT)

rect(s, Inches(0.4), Inches(1.58), Inches(12.5), Inches(0.7),
     fill=C_LGRAY, line=C_ACCENT2, lw=Pt(1))
tb(s, "ℒ(θ) = ℒ_datos(θ)  +  λ_fis · ℒ_fisica(θ)",
   Inches(0.6), Inches(1.68), Inches(12.1), Inches(0.5),
   size=22, bold=True, color=C_BOLD_NUM, align=PP_ALIGN.CENTER)

tb(s, "ℒ_datos: error en condiciones de frontera",
   Inches(0.6), Inches(2.38), Inches(5.5), Inches(0.35), size=16, color=C_BODY)
tb(s, "ℒ_fisica: residuo de Helmholtz en puntos interiores",
   Inches(6.8), Inches(2.38), Inches(6.0), Inches(0.35), size=16, color=C_BODY)

rect(s, Inches(0.4), Inches(2.82), Inches(12.5), Inches(0.025), fill=C_MGRAY)

tb(s, "¿Por que SIREN?  Activacion sinusoidal: φ(z) = sin(ω₀·z)",
   Inches(0.4), Inches(2.92), Inches(12.5), Inches(0.38),
   size=19, bold=True, color=C_ACCENT)

probs_activ = [
    "ReLU / tanh: derivadas de 2.º orden inestables (Helmholtz requiere ∇²E)",
    "Sesgo espectral (spectral bias): aprenden frecuencias bajas primero, fallan para speckle",
    "SIREN: sin(ω₀z) → derivadas de cualquier orden son sinusoides → estabilidad numerica",
    "Inicializacion especial (Sitzmann et al. 2020): distribucion uniforme en [−1, 1] en todas las capas",
    "ω₀ = 1.0 (calibrado para dominio adimensional k = 2π, vs. ω₀ = 30 para imagenes)",
]
bullets(s, probs_activ, Inches(0.4), Inches(3.38), Inches(12.5), Inches(3.5), size=17)

notes(s, "PINNS Y SIREN - 1:45 min\nLas PINNs minimizan la perdida combinando error en frontera + residuo de la EDP.\nSIREN resuelve el problema de derivadas de segundo orden y spectral bias de las activaciones estandar.\nω₀ = 1.0 calibrado para nuestro dominio adimensional.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ARQUITECTURA Y METODOLOGIA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Arquitectura SIREN 5x128 y estrategia de entrenamiento")
slide_num(s, 7)

tb(s, "Arquitectura SIREN 5 x 128",
   Inches(0.4), Inches(1.15), Inches(5.8), Inches(0.38),
   size=19, bold=True, color=C_ACCENT)
arch = [
    "Entrada:  (x, y)  →  2 valores",
    "5 capas ocultas x 128 neuronas",
    "Activacion: sin(ω₀·z),  ω₀ = 1.0",
    "Salida lineal: (E_real, E_imag)  →  2 neuronas",
    "Parametros entrenables: 66,690",
    "Inicializacion: Sitzmann et al. (2020)",
]
bullets(s, arch, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.2), size=18)

for i, (lbl, col) in enumerate([
    ("(x,y)", C_ACCENT2),
    ("128\nsin(ω₀z)", C_ACCENT),
    ("128\nsin(ω₀z) ×3", C_ACCENT),
    ("128\nsin(ω₀z)", C_ACCENT),
    ("(Er,Ei)", C_BOLD_NUM),
]):
    x = Inches(0.45 + i*1.12)
    rect(s, x, Inches(5.2), Inches(1.0), Inches(0.7), fill=col, line=None)
    tb(s, lbl, x+Inches(0.05), Inches(5.26), Inches(0.9), Inches(0.58),
       size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
for xi in [1.6, 2.72, 3.84, 4.96]:
    tb(s, "→", Inches(xi), Inches(5.4), Inches(0.25), Inches(0.3),
       size=14, bold=True, color=C_MGRAY, align=PP_ALIGN.CENTER)

tb(s, "Muestreo y optimizacion",
   Inches(6.7), Inches(1.15), Inches(6.2), Inches(0.38),
   size=19, bold=True, color=C_ACCENT)
meth = [
    "Puntos interiores: LHS 3,000 pts (NB02)",
    "  Cobertura uniforme sin patrones redundantes",
    "Puntos de frontera: 300 pts por borde (x 4 bordes)",
    "Peso de fisica: λ_fis = 0.1",
    "  [Ablacion: 0.01 → 0.1 (optimo) → 1.0 (colapso)]",
    "",
    "Fase 1 - Adam:",
    "  η = 10⁻³,  max 15,000 epocas,  early stopping",
    "Fase 2 - L-BFGS (strong Wolfe):",
    "  Max 1,000 iter,  historial 100",
    "  Ajuste fino de alta precision con curvatura",
]
bullets(s, meth, Inches(6.7), Inches(1.6), Inches(6.2), Inches(4.0), size=16)

notes(s, "ARQUITECTURA - 1:30 min\nSIREN 5x128: 66,690 parametros. Entrada (x,y), salida (E_real, E_imag).\nLHS garantiza cobertura uniforme. Estrategia bifasica: Adam explora globalmente, L-BFGS refina con curvatura.\nλ_fis = 0.1 calibrado: λ=1.0 causa explosion de gradientes en float32.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — NB01 RESULTADOS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB01: Validacion Helmholtz 1D - Comparacion con solucion analitica",
       sub="SIREN 5×64  |  k = 2π  |  E_exacta(x) = cos(kx)  |  1,000 puntos de evaluacion")
slide_num(s, 8)

figure(s, FIG_1D_RES, Inches(0.25), Inches(1.12), Inches(8.55), Inches(5.65),
       cap="Fig. 4.1: PINN vs solucion analitica — Error puntual — Curvas de perdida Adam + L-BFGS")

tb(s, "Metricas clave", Inches(9.0), Inches(1.15), Inches(4.0), Inches(0.38),
   size=18, bold=True, color=C_ACCENT)

metric_box(s, "Error L2", "0.006%", "meta: < 5%",      Inches(9.0),  Inches(1.58))
metric_box(s, "R cuadrado", "1.000000", "ajuste perfecto", Inches(11.35), Inches(1.58))
metric_box(s, "L-BFGS", "202 iter", "de 500 max",      Inches(9.0),  Inches(3.12))
metric_box(s, "Tiempo", "251 s", "total",               Inches(11.35), Inches(3.12))

rect(s, Inches(9.0), Inches(4.65), Inches(4.1), Inches(1.55),
     fill=RGBColor(0xE8, 0xF4, 0xEA), line=C_ACCENT2, lw=Pt(1))
tb(s, "Mejora vs. estado del arte:", Inches(9.1), Inches(4.72),
   Inches(3.9), Inches(0.35), size=14, bold=True, color=C_GREEN)
tb(s, "×415 vs Schoder & Kraxberger\n(2024): error 2.49%",
   Inches(9.1), Inches(5.05), Inches(3.9), Inches(1.0),
   size=20, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

notes(s, "NB01 - 2:00 min\nValidacion contra solucion analitica exacta cos(kx). La curva PINN y la exacta son visualmente indistinguibles.\nError L2 = 0.006%: casi 3 ordenes de magnitud por debajo del umbral. R² = 1.000000. L-BFGS: 202 iter.\nFactor de mejora x415 respecto a Schoder & Kraxberger (2024) con PINNs tanh.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — NB01 METRICAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB01: Analisis detallado de error y convergencia")
slide_num(s, 9)

figure(s, FIG_1D_MET, Inches(0.25), Inches(1.12), Inches(12.8), Inches(5.3),
       cap="Fig. 4.2: Histograma de error absoluto — Error relativo puntual — Correlacion PINN vs Exacta: R2 = 1.000000, Pearson = 1.000000")

tb(s, "Error maximo = 5.77e-5   |   MAE = 3.50e-5   |   "
      "Adam: 15,000 epocas   |   L-BFGS: 202 iteraciones",
   Inches(0.5), Inches(6.52), Inches(12.3), Inches(0.35),
   size=14, italic=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

notes(s, "NB01 METRICAS - 0:45 min\nHistograma: errores del orden de 10^{-5}. Error relativo: media 0.012%. Correlacion PINN vs Exacta: R2 = Pearson = 1.000000.\nConcluimos que la arquitectura y el protocolo son correctos. Procedemos con confianza al caso 2D.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — NB02 RESULTADOS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB02: Helmholtz 2D con campo complejo",
       sub="SIREN 5×128  |  k = 2π  |  LHS 3,000 pts  |  E_exacta = e^(i(kₓx + k_yy))")
slide_num(s, 10)

figure(s, FIG_2D_RES, Inches(0.2), Inches(1.12), Inches(8.7), Inches(5.7),
       cap="Fig. 4.3: E_real y E_imag — solucion exacta (fila sup.) vs prediccion PINN (fila inf.) — Error puntual y curvas de perdida")

tb(s, "Metricas NB02", Inches(9.1), Inches(1.15), Inches(4.0), Inches(0.38),
   size=18, bold=True, color=C_ACCENT)

rows = [
    ("Metrica",          "E_real",   "E_imag"),
    ("Error L2 (%)",     "0.214",    "0.127"),
    ("R cuadrado",       "0.999995", "0.999998"),
    ("Pearson",          "0.999998", "0.999999"),
    ("Error L2 prom.",   "0.171%",   ""),
    ("Epocas Adam",      "8,737",    ""),
    ("Iter. L-BFGS",     "1,035",    ""),
    ("Tiempo total",     "299 s",    ""),
]
t_top = 1.60
for i, (m, v1, v2) in enumerate(rows):
    is_hdr  = (i == 0)
    is_prom = (i == 4)
    bg_r = C_ACCENT if is_hdr else (RGBColor(0xE8, 0xF4, 0xEA) if is_prom else (C_LGRAY if i % 2 == 0 else C_WHITE))
    tc   = C_WHITE if is_hdr else (C_GREEN if is_prom else C_BODY)
    rect(s, Inches(9.1), Inches(t_top+i*0.5), Inches(4.0), Inches(0.48),
         fill=bg_r, line=C_ACCENT2, lw=Pt(0.5))
    tb(s, m,  Inches(9.15), Inches(t_top+i*0.5+0.07), Inches(1.7), Inches(0.35),
       size=12, bold=(is_hdr or is_prom), color=tc)
    tb(s, v1, Inches(10.9), Inches(t_top+i*0.5+0.07), Inches(1.1), Inches(0.35),
       size=12, bold=(is_hdr or is_prom),
       color=C_BOLD_NUM if is_prom else tc, align=PP_ALIGN.CENTER)
    tb(s, v2, Inches(12.05), Inches(t_top+i*0.5+0.07), Inches(0.95), Inches(0.35),
       size=12, bold=is_hdr, color=tc, align=PP_ALIGN.CENTER)

notes(s, "NB02 - 1:45 min\nMapas de calor: E_real y E_imag. Solucion exacta (arriba) vs PINN (abajo): indistinguibles visualmente.\nError L2 promedio 0.171%: mas de un orden por debajo del umbral. R² > 0.9999. Tiempo: 299 s.\nFactor de mejora x11.2 vs Schoder & Kraxberger. Ampliar de d=64 a d=128 fue determinante: 0.436% → 0.171%.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — NB02 ROBUSTEZ MULTI-SEMILLA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "NB02: Analisis de error y reproducibilidad multi-semilla")
slide_num(s, 11)

figure(s, FIG_2D_MET, Inches(0.2), Inches(1.12), Inches(7.6), Inches(5.15),
       cap="Fig. 4.4: Distribucion de errores E_real / E_imag — Mapa de error puntual — Correlacion PINN vs Exacta")

tb(s, "Robustez multi-semilla (NB02, λ_fis = 0.1)",
   Inches(8.0), Inches(1.15), Inches(5.0), Inches(0.38),
   size=16, bold=True, color=C_ACCENT)

seeds = [
    ("SEED",       "Error L2",       "Epocas",  "Tiempo"),
    ("42",         "0.171%",         "8,737",   "299 s"),
    ("123",        "0.126%",         "8,012",   "448 s"),
    ("777",        "0.168%",         "13,560",  "346 s"),
    ("Media ± DE", "0.155 ± 0.020%", "--",      "--"),
]
t_top = 1.60
for i, row in enumerate(seeds):
    is_hdr = (i == 0)
    is_sum = (i == 4)
    bg_r = C_ACCENT if is_hdr else (RGBColor(0xE8, 0xF4, 0xEA) if is_sum else (C_LGRAY if i % 2 == 0 else C_WHITE))
    tc   = C_WHITE if is_hdr else (C_GREEN if is_sum else C_BODY)
    rect(s, Inches(8.0), Inches(t_top+i*0.55), Inches(5.0), Inches(0.52),
         fill=bg_r, line=C_ACCENT2, lw=Pt(0.5))
    col_x = [8.05, 9.3, 10.4, 11.5]
    col_w = [1.22, 1.07, 1.07, 1.42]
    for j, (val, cw) in enumerate(zip(row, col_w)):
        tb(s, val, Inches(col_x[j]), Inches(t_top+i*0.55+0.09),
           Inches(cw), Inches(0.35), size=13, bold=(is_hdr or is_sum),
           color=tc, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

rect(s, Inches(8.0), Inches(4.65), Inches(5.0), Inches(1.3),
     fill=RGBColor(0xE8, 0xF4, 0xEA), line=C_ACCENT2, lw=Pt(1))
tb(s, "Varianza inter-semilla: 13% relativo",
   Inches(8.15), Inches(4.73), Inches(4.7), Inches(0.38),
   size=16, bold=True, color=C_GREEN)
tb(s, "El metodo es estadisticamente robusto:\nconverge independientemente de la inicializacion.",
   Inches(8.15), Inches(5.1), Inches(4.7), Inches(0.75),
   size=15, color=C_BODY)

notes(s, "NB02 ROBUSTEZ - 1:00 min\nRepetimos con seeds 42, 123, 777. Error entre 0.126% y 0.171%. Media 0.155%, DE 0.020%.\nVarianza inter-semilla 13% relativo: el metodo converge robustamente sin importar la inicializacion.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — COMPARATIVA ESTADO DEL ARTE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Comparativa con el estado del arte")
slide_num(s, 12)

tb(s, "Referencia: Schoder & Kraxberger (2024) - PINNs convencionales (activacion tanh) para Helmholtz 3D",
   Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.35),
   size=14, italic=True, color=C_ACCENT2)

hdrs = ["Metodo", "Dimension", "Error L2 (%)", "Factor de mejora"]
rows_c = [
    ("Schoder & Kraxberger (2024)", "1D", "2.490", "--"),
    ("Schoder & Kraxberger (2024)", "2D", "1.910", "--"),
    ("PINN-SIREN (este trabajo)",   "1D", "0.006", "×415"),
    ("PINN-SIREN (este trabajo)",   "2D", "0.171", "×11.2"),
]
cw = [4.5, 1.5, 1.8, 2.4]
cx = [0.4, 4.9, 6.4, 8.2]
t_top = 1.6

for j, (h, w) in enumerate(zip(hdrs, cw)):
    rect(s, Inches(cx[j]), Inches(t_top), Inches(w), Inches(0.52),
         fill=C_ACCENT, line=None)
    tb(s, h, Inches(cx[j]+0.05), Inches(t_top+0.08),
       Inches(w-0.1), Inches(0.35), size=14, bold=True,
       color=C_WHITE, align=PP_ALIGN.CENTER)

for i, (m, d, e, f) in enumerate(rows_c):
    rt = t_top + (i+1)*0.62
    is_ours = (i >= 2)
    bg_r = RGBColor(0xE8, 0xF4, 0xEA) if is_ours else (C_LGRAY if i % 2 == 0 else C_WHITE)
    for j, (val, w) in enumerate(zip([m, d, e, f], cw)):
        rect(s, Inches(cx[j]), Inches(rt), Inches(w), Inches(0.6),
             fill=bg_r, line=C_ACCENT2, lw=Pt(0.5))
        tc = C_GREEN if (is_ours and j in [2, 3]) else C_BODY
        tb(s, val, Inches(cx[j]+0.05), Inches(rt+0.1),
           Inches(w-0.1), Inches(0.42), size=14, bold=is_ours, color=tc,
           align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(4.45), Inches(5.5), Inches(1.65),
     fill=C_LGRAY, line=C_ACCENT, lw=Pt(1.5))
tb(s, "×415 en 1D",
   Inches(0.5), Inches(4.55), Inches(5.3), Inches(0.65),
   size=36, bold=True, color=C_BOLD_NUM, align=PP_ALIGN.CENTER)
tb(s, "SIREN vs PINN convencional — Helmholtz 1D",
   Inches(0.5), Inches(5.18), Inches(5.3), Inches(0.75),
   size=14, color=C_BODY, align=PP_ALIGN.CENTER)

rect(s, Inches(6.4), Inches(4.45), Inches(5.5), Inches(1.65),
     fill=C_LGRAY, line=C_ACCENT, lw=Pt(1.5))
tb(s, "×11.2 en 2D",
   Inches(6.5), Inches(4.55), Inches(5.3), Inches(0.65),
   size=36, bold=True, color=C_BOLD_NUM, align=PP_ALIGN.CENTER)
tb(s, "SIREN 2D vs PINN convencional — Helmholtz 2D/3D",
   Inches(6.5), Inches(5.18), Inches(5.3), Inches(0.75),
   size=14, color=C_BODY, align=PP_ALIGN.CENTER)

tb(s, "Nota: comparacion indicativa. Schoder & Kraxberger operan en 3D con dominio fisico y activacion tanh; "
      "este trabajo en 1D/2D adimensional con SIREN.",
   Inches(0.4), Inches(6.2), Inches(11.8), Inches(0.35),
   size=11, italic=True, color=C_ACCENT2)

notes(s, "COMPARATIVA - 1:15 min\n1D: su error 2.49%, el nuestro 0.006%. Factor x415.\n2D: su error 1.91%, el nuestro 0.171%. Factor x11.2.\nComparacion indicativa: ellos en 3D con dominio fisico, nosotros en 1D/2D adimensional con SIREN.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DISCUSION (sin NB03)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Discusion: Hallazgos principales del avance")
slide_num(s, 13)

boxes_d = [
    ("NB01\nHelmholtz 1D", "L² = 0.006%\nR² = 1.000000\n251 s"),
    ("NB02\nHelmholtz 2D", "L² = 0.171%\nR² > 0.9999\n299 s"),
]
for i, (tit, met) in enumerate(boxes_d):
    rect(s, Inches(1.0+i*5.8), Inches(1.2), Inches(5.3), Inches(1.6),
         fill=C_ACCENT, line=None)
    tb(s, tit, Inches(1.1+i*5.8), Inches(1.28),
       Inches(5.1), Inches(0.58), size=16, bold=True,
       color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(s, met, Inches(1.1+i*5.8), Inches(1.86),
       Inches(5.1), Inches(0.82), size=16,
       color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# Caja NB03 pendiente
rect(s, Inches(8.6), Inches(1.2), Inches(4.3), Inches(1.6),
     fill=C_AMBER, line=None)
tb(s, "NB03\nSpeckle optico", Inches(8.7), Inches(1.28),
   Inches(4.1), Inches(0.58), size=16, bold=True,
   color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "PENDIENTE\nProximo experimento", Inches(8.7), Inches(1.86),
   Inches(4.1), Inches(0.82), size=15,
   color=C_WHITE, align=PP_ALIGN.CENTER)

tb(s, "Hallazgos clave:",
   Inches(0.4), Inches(3.0), Inches(12.5), Inches(0.38),
   size=19, bold=True, color=C_ACCENT)

hall = [
    "Activacion sinusoidal SIREN supera ×415 a PINNs convencionales (tanh) para ecuaciones de onda",
    "Optimizacion hibrida Adam + L-BFGS es esencial: Adam evita minimos locales, L-BFGS refina con curvatura",
    "ω₀ = 1.0 y λ_fis = 0.1 son configuraciones no triviales calibradas empiricamente",
    "Robustez multi-semilla confirmada: varianza inter-semilla 13% relativo con seeds {42, 123, 777}",
    "La misma arquitectura SIREN 5×128 es la base para NB03 (speckle) sin modificacion de disenho",
]
bullets(s, hall, Inches(0.4), Inches(3.45), Inches(12.5), Inches(3.8), size=17)

notes(s, "DISCUSION - 1:00 min\nNB01 y NB02 completados con amplio margen. NB03 es el siguiente paso.\nHallazgos clave: SIREN supera x415 a tanh; la combinacion Adam + L-BFGS es esencial; calibracion de hiperparametros documentada; robustez multi-semilla confirmada.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — TRABAJO FUTURO (NB03 como primero)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Trabajo futuro y proximos pasos")
slide_num(s, 14)

tb(s, "Experimentos planificados (orden de prioridad):",
   Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.38),
   size=20, bold=True, color=C_ACCENT)

future = [
    ("NB03", "Simulacion de speckle optico  [PRIORIDAD]",
     "Condicion de frontera: φ(x)~U(0,2π) en y=0, N_φ=256 pts. Validar con criterio de Goodman: |C-1| < 0.1.",
     C_AMBER),
    ("NB04", "Benchmark vs FEM (FEniCSx)",
     "Comparar PINN-SIREN vs Metodo de Elementos Finitos. Metrica: S = T_FEM / T_PINN. Cuantificar aceleracion.",
     C_ACCENT2),
    ("NB05", "Medios inhomogeneos (GRIN)",
     "Extender a: ∇²E + k²(r)E = 0 con indice variable. Aplicacion: fibras opticas con gradiente de indice.",
     C_ACCENT2),
    ("NB06", "Condicion de radiacion de Sommerfeld",
     "Agregar: ∂E/∂n − ikE = 0 en y=1 para eliminar reflexiones espurias y mejorar validacion estadistica.",
     C_ACCENT2),
]
top = 1.65
for code, title, desc, col in future:
    rect(s, Inches(0.4), Inches(top), Inches(0.9), Inches(1.0), fill=col, line=None)
    tb(s, code, Inches(0.4), Inches(top+0.05),
       Inches(0.9), Inches(0.45), size=13, bold=True,
       color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, Inches(1.45), Inches(top+0.04),
       Inches(11.3), Inches(0.38), size=17, bold=True, color=C_ACCENT)
    tb(s, desc, Inches(1.45), Inches(top+0.42),
       Inches(11.3), Inches(0.5), size=15, color=C_BODY)
    top += 1.1

notes(s, "TRABAJO FUTURO - 0:45 min\nNB03 es la prioridad inmediata: aplicar el modelo validado a la generacion de speckle con fase aleatoria.\nNB04: benchmark FEM para cuantificar S = T_FEM / T_PINN.\nNB05/NB06: extensiones a medios inhomogeneos y condicion de Sommerfeld.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — RESUMEN (sin columna speckle)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s)
header(s, "Resumen: Que se logro en este primer avance")
slide_num(s, 15)

cols_r = [
    ("Precision demostrada",
     ["Error L2 = 0.006% (1D)", "Error L2 = 0.171% (2D)",
      "Meta < 5% superada", "R² = 1.000000 en 1D"]),
    ("Robustez confirmada",
     ["Multi-semilla {42, 123, 777}", "Varianza 13% relativo",
      "Convergencia independiente", "de inicializacion"]),
    ("Infraestructura solida",
     ["Codigo modular en src/", "2 notebooks reproducibles",
      "SEED=42 garantiza reprod.", "GPU RTX 5050, CUDA 12.6"]),
]
for i, (tit, items) in enumerate(cols_r):
    x = Inches(0.35 + i*4.3)
    rect(s, x, Inches(1.2), Inches(4.1), Inches(0.55), fill=C_ACCENT, line=None)
    tb(s, tit, x+Inches(0.08), Inches(1.25), Inches(3.94), Inches(0.45),
       size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        bg_it = C_LGRAY if j % 2 == 0 else C_WHITE
        rect(s, x, Inches(1.75+j*0.75), Inches(4.1), Inches(0.72),
             fill=bg_it, line=C_ACCENT2, lw=Pt(0.5))
        tb(s, item, x+Inches(0.08), Inches(1.8+j*0.75),
           Inches(3.94), Inches(0.6), size=15, color=C_BODY)

rect(s, Inches(0.35), Inches(4.85), Inches(12.6), Inches(0.7),
     fill=C_LGRAY, line=C_ACCENT, lw=Pt(1.5))
tb(s, "NB01 y NB02 completos — L2 muy por debajo del umbral del 5%  ✓"
      "   |   NB03 (speckle) y NB04 (benchmark FEM): proximos pasos",
   Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.5),
   size=15, italic=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

notes(s, "RESUMEN - 0:20 min\nNB01 y NB02 completos: precision muy por debajo del umbral. Robustez multi-semilla confirmada. Infraestructura reproducible lista.\nNB03 y NB04 son el trabajo futuro inmediato.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CIERRE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank())
bg(s, C_ACCENT)
rect(s, 0, 0, SW, Inches(0.12), fill=C_WHITE)
rect(s, 0, SH-Inches(0.12), SW, Inches(0.12), fill=C_WHITE)

tb(s, "Gracias por su atencion",
   Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.85),
   size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(1.8), Inches(2.0), Inches(9.7), Inches(0.03),
     fill=RGBColor(0xCC, 0xDD, 0xEE))

tb(s, "Roberto Hernandez Estrada",
   Inches(0.5), Inches(2.15), Inches(12.3), Inches(0.48),
   size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(s, "robertohernandezestrd@gmail.com   |   Matricula: 252H21004",
   Inches(0.5), Inches(2.62), Inches(12.3), Inches(0.35),
   size=14, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(1.8), Inches(3.1), Inches(9.7), Inches(0.03),
     fill=RGBColor(0x80, 0xA0, 0xBB))

tb(s, "Agradecimientos",
   Inches(0.5), Inches(3.28), Inches(12.3), Inches(0.38),
   size=18, bold=True, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

acks = [
    "Dr. Jose Adan Hernandez Nolasco  -  Director de tesis",
    "Dr. Oscar Alberto Chavez Bosquez  -  Codirector de tesis",
    "DACYTI-UJAT  -  Programa de Maestria en Ciencias de la Computacion",
]
for k, ack in enumerate(acks):
    tb(s, ack, Inches(1.5), Inches(3.75+k*0.5), Inches(10.3), Inches(0.42),
       size=17, color=C_WHITE, align=PP_ALIGN.CENTER)

tb(s, "Python 3.11  |  PyTorch 2.4  |  CUDA 12.6  |  GPU NVIDIA RTX 5050",
   Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.38),
   size=13, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

rect(s, Inches(1.8), Inches(6.1), Inches(9.7), Inches(0.03),
     fill=RGBColor(0x80, 0xA0, 0xBB))

tb(s, "Preguntas?",
   Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
   size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

notes(s, "CIERRE - 0:10 min\nMuchas gracias. Quedo a sus ordenes para preguntas.")

# ── Guardar ────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)
print(f"Presentacion guardada: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
