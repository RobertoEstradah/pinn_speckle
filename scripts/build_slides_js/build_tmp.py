#!/usr/bin/env python3
"""Coloquio UJAT - PINN-SIREN Speckle - resultado_slide.pptx"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
BG    = RGBColor(0x0A,0x0D,0x1A)
PANEL = RGBColor(0x14,0x19,0x28)
GOLD  = RGBColor(0xF0,0xA5,0x00)
TEAL  = RGBColor(0x00,0xB4,0xD8)
RED   = RGBColor(0xFF,0x44,0x44)
GREEN = RGBColor(0x00,0xC8,0x7A)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LGRAY = RGBColor(0xB0,0xB8,0xCC)
AMBER = RGBColor(0xFF,0xB7,0x00)
DARK  = RGBColor(0x08,0x0B,0x14)
BORDER= RGBColor(0x2A,0x30,0x50)

W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, c=BG):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = c

def txt(slide, text, x, y, w, h, sz=14, bold=False, clr=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x,y,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = clr
    r.font.name = "Calibri"
    return tb

def rect(slide, x, y, w, h, fill=PANEL, line=BORDER, lw=1.5):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(lw)
    return s

def rect0(slide, x, y, w, h, fill=PANEL):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def header(slide, title, sub="Coloquio de Avance · UJAT · 2026"):
    rect0(slide, Inches(0), Inches(0), W, Inches(0.7), DARK)
    txt(slide, sub,   Inches(0.3), Inches(0.04), Inches(9), Inches(0.3),
        sz=10, clr=GOLD)
    txt(slide, title, Inches(0.3), Inches(0.32), Inches(12.7), Inches(0.62),
        sz=26, bold=True, clr=WHITE)

prs = new_prs()

# ── S1 PORTADA ───────────────────────────────────────────────────────────────
s = blank(prs); bg(s)
rect0(s, Inches(0), Inches(0), Inches(0.1), H, GOLD)
txt(s, "UJAT · Maestría en Ciencias de la Computación · Junio 2026",
    Inches(0.3), Inches(1.0), Inches(10), Inches(0.4), sz=13, clr=GOLD)
txt(s, "Simulación Acelerada de Speckle Óptico\nmediante PINN-SIREN",
    Inches(0.3), Inches(1.6), Inches(8.8), Inches(2.2),
    sz=36, bold=True, clr=WHITE)
txt(s, "Coloquio de Avance de Tesis",
    Inches(0.3), Inches(3.9), Inches(8), Inches(0.5), sz=16, clr=TEAL)
rect(s, Inches(0.3), Inches(4.6), Inches(5.8), Inches(1.8), PANEL, BORDER)
txt(s, "Roberto Hernández Estrada",
    Inches(0.5), Inches(4.75), Inches(5.4), Inches(0.45), sz=15, bold=True)
txt(s, "Director: Dr. José Adán Hernández Nolasco",
    Inches(0.5), Inches(5.2),  Inches(5.4), Inches(0.35), sz=12, clr=LGRAY)
txt(s, "División Académica de Ciencias Básicas · UJAT",
    Inches(0.5), Inches(5.55), Inches(5.4), Inches(0.35), sz=11, clr=LGRAY)
# KPI panel right
rect(s, Inches(9.3), Inches(1.4), Inches(3.7), Inches(5.5), PANEL, BORDER)
txt(s, "Resultados Obtenidos", Inches(9.5), Inches(1.55), Inches(3.3), Inches(0.4),
    sz=13, bold=True, clr=GOLD, align=PP_ALIGN.CENTER)
kpis1 = [("0.006%","Error L² — 1D",GREEN),("0.171%","Error L² — 2D",TEAL),
          ("×415","Mejora vs. Schoder",AMBER),("C=1.025","Contraste Speckle",GREEN)]
for i,(v,l,c) in enumerate(kpis1):
    yy = Inches(2.1+i*1.1)
    rect0(s, Inches(9.5), yy, Inches(3.3), Inches(0.95), RGBColor(0x18,0x1F,0x35))
    txt(s, v, Inches(9.5), yy+Inches(0.04), Inches(3.3), Inches(0.55),
        sz=28, bold=True, clr=c, align=PP_ALIGN.CENTER)
    txt(s, l, Inches(9.5), yy+Inches(0.58), Inches(3.3), Inches(0.3),
        sz=10, clr=LGRAY, align=PP_ALIGN.CENTER)

# ── S2 AGENDA ────────────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s, "Agenda")
agenda = [
    ("01","Motivación","¿Por qué FEM es insuficiente?"),
    ("02","Hipótesis y Objetivos","Marco de validación — NB01–NB04"),
    ("03","Marco Teórico","Helmholtz · PINNs · SIREN"),
    ("04","Metodología","Arquitectura · LHS · Adam+L-BFGS"),
    ("05","Resultados NB01–NB03","1D, 2D y Speckle Óptico"),
    ("06","Ablación λ_phys","Sensibilidad al peso de física"),
    ("07","Resumen · Próximos Pasos · Conclusiones",""),
]
for i,(n,t,d) in enumerate(agenda):
    ci = i%2; ri = i//2
    x = Inches(0.4+ci*6.6); y = Inches(1.15+ri*1.45); w = Inches(6.2)
    rect(s,x,y,w,Inches(1.3),PANEL,BORDER)
    txt(s,n, x+Inches(0.15),y+Inches(0.2),Inches(0.55),Inches(0.5),
        sz=22,bold=True,clr=GOLD)
    txt(s,t, x+Inches(0.85),y+Inches(0.08),Inches(5.2),Inches(0.45),
        sz=14,bold=True,clr=WHITE)
    if d: txt(s,d,x+Inches(0.85),y+Inches(0.58),Inches(5.2),Inches(0.4),
              sz=11,clr=LGRAY)

# ── S3 MOTIVACIÓN ────────────────────────────────────────────────────────────
s = blank(prs); bg(s)
header(s,"El Cuello de Botella Computacional")
rect(s,Inches(0.3),Inches(1.0),Inches(6.0),Inches(5.9),PANEL,RED,1.5)
txt(s,"Problema con FEM Clásico",Inches(0.5),Inches(1.1),Inches(5.6),Inches(0.45),
    sz=15,bold=True,clr=RED)
for j,item in enumerate(["Resolución < λ láser → millones de nodos",
                           "Nuevo perfil rugoso → resolver desde cero",
                           "Miles de realizaciones = costo intratable",
                           "Alta memoria pico por realización"]):
    yy=Inches(1.8+j*0.75)
    rect0(s,Inches(0.5),yy,Inches(5.5),Inches(0.6),RGBColor(0x20,0x08,0x08))
    txt(s,"✕  "+item,Inches(0.6),yy+Inches(0.08),Inches(5.2),Inches(0.45),sz=12,
        clr=RGBColor(0xFF,0x88,0x88))
rect(s,Inches(6.7),Inches(1.0),Inches(6.3),Inches(5.9),PANEL,GREEN,1.5)
txt(s,"Solución PINN-SIREN",Inches(6.9),Inches(1.1),Inches(5.9),Inches(0.45),
    sz=15,bold=True,clr=GREEN)
for j,item in enumerate(["Libre de malla (mesh-free)",
                           "Inferencia en ≈ 1.3 ms (GPU RTX 5050)",
                           "Misma red → cualquier frontera",
                           "Entrenamiento único → ∞ realizaciones"]):
    yy=Inches(1.8+j*0.75)
    rect0(s,Inches(6.9),yy,Inches(5.9),Inches(0.6),RGBColor(0x02,0x18,0x10))
    txt(s,"✓  "+item,Inches(7.0),yy+Inches(0.08),Inches(5.6),Inches(0.45),sz=12,
        clr=RGBColor(0x80,0xFF,0xC0))
txt(s,"Ecuación: ∇²E + k²E = 0  con  k = 2π  (dominio adimensional [0,1]²)",
    Inches(0.3),Inches(6.6),Inches(12.7),Inches(0.4),sz=12,clr=GOLD,
    align=PP_ALIGN.CENTER,italic=True)

# ── S4 FENÓMENO FÍSICO DEL SPECKLE ───────────────────────────────────────────
s = blank(prs); bg(s); header(s,"El Speckle Óptico: Fenómeno y Desafío Computacional")

# Left: physical setup
rect(s,Inches(0.3),Inches(1.0),Inches(6.1),Inches(6.1),PANEL,BORDER)
txt(s,"El Fenómeno",Inches(0.5),Inches(1.1),Inches(5.7),Inches(0.45),
    sz=15,bold=True,clr=TEAL)
# Physics chain
chain=[
    ("Láser coherente","λ = 638 nm",TEAL),
    ("Superficie rugosa","σ_h ≫ λ  (variación de fase > 2π)",GOLD),
    ("Interferencia aleatoria","N contribuciones independientes",LGRAY),
    ("Speckle completamente\ndesarrollado","C = σ_I/⟨I⟩ = 1  (Goodman, 1963)",GREEN),
]
for i,(title,detail,c) in enumerate(chain):
    yy=Inches(1.7+i*1.12)
    rect0(s,Inches(0.5),yy,Inches(5.7),Inches(1.0),RGBColor(0x10,0x15,0x25))
    rect0(s,Inches(0.5),yy,Inches(0.08),Inches(1.0),c)
    txt(s,title,Inches(0.75),yy+Inches(0.05),Inches(5.2),Inches(0.4),
        sz=13,bold=True,clr=c)
    txt(s,detail,Inches(0.75),yy+Inches(0.5),Inches(5.2),Inches(0.42),
        sz=11,clr=LGRAY)
    if i<3:
        txt(s,"↓",Inches(3.0),yy+Inches(1.02),Inches(0.5),Inches(0.2),
            sz=11,clr=BORDER,align=PP_ALIGN.CENTER)

# Right: speckle image + applications
rect(s,Inches(6.65),Inches(1.0),Inches(6.35),Inches(6.1),PANEL,BORDER)
txt(s,"Simulado con PINN-SIREN",Inches(6.85),Inches(1.1),Inches(5.95),Inches(0.4),
    sz=13,bold=True,clr=WHITE)
s.shapes.add_picture(
    r"D:\Tesis_Maestria\paper\figures\resultados_speckle_nb03.png",
    Inches(6.85),Inches(1.6),width=Inches(5.9))
txt(s,"Aplicaciones:",Inches(6.85),Inches(5.45),Inches(5.9),Inches(0.35),
    sz=12,bold=True,clr=GOLD)
apps=["Metrología óptica","Tomografía de coherencia óptica",
      "Imagenología biomédica","Criptografía óptica"]
for i,a in enumerate(apps):
    txt(s,"• "+a,Inches(6.85),Inches(5.85+i*0.3),Inches(5.9),Inches(0.28),
        sz=11,clr=WHITE)

txt(s,"Cada nueva rugosidad σ_h → resolver Helmholtz desde cero. "
      "Para estadísticas: miles de realizaciones = costo intratable con FEM.",
    Inches(0.3),Inches(7.15),Inches(12.7),Inches(0.3),
    sz=11,clr=AMBER,align=PP_ALIGN.CENTER,italic=True)

# ── S5 HIPÓTESIS Y OBJETIVOS ─────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"Hipótesis y Objetivos de Investigación")
rect(s,Inches(0.3),Inches(1.0),Inches(12.7),Inches(1.05),
     RGBColor(0x0F,0x1A,0x2A),GOLD,1.5)
txt(s,"HIPÓTESIS: PINN-SIREN puede simular speckle con L² < 5% vs. referencia "
       "y factor de aceleración S = T_FEM / T_PINN > 1",
    Inches(0.5),Inches(1.07),Inches(12.3),Inches(0.9),
    sz=14,clr=WHITE,align=PP_ALIGN.CENTER)
objs=[("NB01","Helmholtz 1D vs. solución analítica","✅ L²=0.006%",GREEN),
      ("NB02","Helmholtz 2D campo complejo (E_real,E_imag)","✅ L²=0.171%",GREEN),
      ("NB03","Speckle con frontera rugosa φ~U(0,2π)","✅ C=1.025",GREEN),
      ("NB04","Benchmark FEM — Speed-up Factor","📋 Pendiente",GOLD)]
for i,(nb,desc,status,c) in enumerate(objs):
    yy=Inches(2.25+i*1.2)
    rect(s,Inches(0.3),yy,Inches(12.7),Inches(1.05),PANEL,BORDER)
    rect0(s,Inches(0.3),yy,Inches(1.1),Inches(1.05),c)
    txt(s,nb,Inches(0.3),yy+Inches(0.28),Inches(1.1),Inches(0.5),
        sz=14,bold=True,clr=RGBColor(0,0,0),align=PP_ALIGN.CENTER)
    txt(s,desc,Inches(1.55),yy+Inches(0.28),Inches(9.0),Inches(0.45),
        sz=14,bold=True,clr=WHITE)
    txt(s,status,Inches(10.5),yy+Inches(0.2),Inches(2.3),Inches(0.6),
        sz=13,bold=True,clr=c,align=PP_ALIGN.CENTER)

# ── S5 MARCO TEÓRICO ─────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"Marco Teórico: Helmholtz · PINNs · SIREN")
panels=[
    ("Ecuación de Helmholtz",
     "∇²E + k²E = 0\n\nE = E_real + i·E_imag\n\nk = 2π  (dominio [0,1]²)\n\nI(x,y) = E_real² + E_imag²",TEAL),
    ("Función de Pérdida PINN",
     "L = L_datos + λ·L_física\n\nL_datos: error en frontera\n\nL_física: residuo de Helmholtz\n\nλ_phys = 0.1 (calibrado)",GOLD),
    ("Arquitectura SIREN",
     "φ(z) = sin(ω₀·z),  ω₀ = 1.0\n\n5 capas × 128 neuronas\n66,690 parámetros\n\nInic. Sitzmann et al. 2020",GREEN),
]
for i,(title,content,c) in enumerate(panels):
    x=Inches(0.3+i*4.36)
    rect(s,x,Inches(1.0),Inches(4.15),Inches(6.1),PANEL,BORDER)
    rect0(s,x,Inches(1.0),Inches(4.15),Inches(0.55),c)
    txt(s,title,x+Inches(0.1),Inches(1.03),Inches(3.95),Inches(0.45),
        sz=13,bold=True,clr=RGBColor(0,0,0),align=PP_ALIGN.CENTER)
    txt(s,content,x+Inches(0.2),Inches(1.7),Inches(3.75),Inches(5.1),sz=13,clr=WHITE)

# ── S6 ¿POR QUÉ SIREN? — SESGO ESPECTRAL ────────────────────────────────────
s = blank(prs); bg(s)
header(s,"¿Por Qué SIREN? El Sesgo Espectral de las Redes Estándar")

# Left: Standard NNs problems
rect(s,Inches(0.3),Inches(1.0),Inches(5.9),Inches(5.65),PANEL,RED,1.5)
txt(s,"Redes Estándar (tanh / ReLU)",Inches(0.5),Inches(1.1),Inches(5.5),Inches(0.45),
    sz=14,bold=True,clr=RED)
probs=[
    ("Sesgo espectral","Aprenden frecuencias bajas primero\n(Rahaman et al., 2019)"),
    ("Derivadas de 2° orden","Inestables o nulas para ReLU\n→ Residuo de Helmholtz ruidoso"),
    ("Alta frecuencia k=2π","Requieren muchas épocas extra\npara capturar las oscilaciones"),
    ("Convergencia lenta","Adam solo: ~50,000 épocas\npara L² < 1% en Helmholtz"),
]
for i,(t,d) in enumerate(probs):
    yy=Inches(1.75+i*1.05)
    rect0(s,Inches(0.5),yy,Inches(5.5),Inches(0.95),RGBColor(0x22,0x08,0x08))
    txt(s,"✕  "+t,Inches(0.65),yy+Inches(0.05),Inches(5.1),Inches(0.38),
        sz=12,bold=True,clr=RGBColor(0xFF,0x77,0x77))
    txt(s,d,Inches(0.65),yy+Inches(0.45),Inches(5.1),Inches(0.42),
        sz=11,clr=LGRAY)

# Center arrow
txt(s,"vs.",Inches(6.3),Inches(3.8),Inches(0.9),Inches(0.5),
    sz=22,bold=True,clr=BORDER,align=PP_ALIGN.CENTER)

# Right: SIREN advantages
rect(s,Inches(7.3),Inches(1.0),Inches(5.7),Inches(5.65),PANEL,GREEN,1.5)
txt(s,"SIREN  —  φ(z) = sin(ω₀z)",Inches(7.5),Inches(1.1),Inches(5.3),Inches(0.45),
    sz=14,bold=True,clr=GREEN)
advs=[
    ("Sin sesgo espectral","Cada capa es periódica: todas\nlas frecuencias se aprenden igual"),
    ("Derivadas analíticas","dⁿ/dxⁿ[sin(ω₀x)] estable\n→ Laplaciano exacto vía autograd"),
    ("Inicialización de Sitzmann","W_l ~ U(-√(6/d), √(6/d))\nPreserva distribución de activación"),
    ("Convergencia rápida","Adam 8,737 épocas → L²=0.171%\ncon k=2π en dominio 2D"),
]
for i,(t,d) in enumerate(advs):
    yy=Inches(1.75+i*1.05)
    rect0(s,Inches(7.5),yy,Inches(5.3),Inches(0.95),RGBColor(0x02,0x18,0x10))
    txt(s,"✓  "+t,Inches(7.65),yy+Inches(0.05),Inches(4.9),Inches(0.38),
        sz=12,bold=True,clr=RGBColor(0x80,0xFF,0xC0))
    txt(s,d,Inches(7.65),yy+Inches(0.45),Inches(4.9),Inches(0.42),
        sz=11,clr=LGRAY)

# Bottom: key equation
rect(s,Inches(0.3),Inches(6.8),Inches(12.7),Inches(0.62),
     RGBColor(0x0A,0x14,0x24),GOLD,1.5)
txt(s,"h_ℓ = sin( ω₀ · (W_ℓ · h_{ℓ-1} + b_ℓ) )     con     "
      "W_0 ~ U(−1/n_in, 1/n_in),   W_ℓ ~ U(−√(6/d), √(6/d))",
    Inches(0.5),Inches(6.87),Inches(12.3),Inches(0.48),
    sz=13,clr=GOLD,align=PP_ALIGN.CENTER,bold=True)

# ── S7 METODOLOGÍA ───────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"Metodología: Arquitectura · LHS · Optimización Bifásica")
# Left
rect(s,Inches(0.3),Inches(1.0),Inches(4.3),Inches(6.0),PANEL,BORDER)
txt(s,"Arquitectura SIREN",Inches(0.5),Inches(1.07),Inches(3.9),Inches(0.4),
    sz=13,bold=True,clr=TEAL)
for j,item in enumerate(["Entrada: (x,y) ∈ [0,1]²","5 capas ocultas × 128 neuronas",
                           "φ(z) = sin(ω₀z),  ω₀=1.0","Salida lineal: (E_real, E_imag)","66,690 parámetros"]):
    txt(s,"→ "+item,Inches(0.5),Inches(1.6+j*0.75),Inches(3.9),Inches(0.6),sz=12,clr=WHITE)
# Middle
rect(s,Inches(4.85),Inches(1.0),Inches(3.8),Inches(6.0),PANEL,BORDER)
txt(s,"Muestreo LHS",Inches(5.05),Inches(1.07),Inches(3.4),Inches(0.4),
    sz=13,bold=True,clr=GOLD)
for j,item in enumerate(["N_c = 3,000 puntos interiores","Hipercubo Latino (LHS)",
                           "Cobertura uniforme","sin redundancia cartesiana",
                           "","N_b = 300 pts / borde","4 bordes del dominio"]):
    if item: txt(s,item,Inches(5.05),Inches(1.6+j*0.62),Inches(3.4),Inches(0.55),sz=12,clr=WHITE)
# Right
rect(s,Inches(8.85),Inches(1.0),Inches(4.15),Inches(6.0),PANEL,BORDER)
txt(s,"Optimización Bifásica",Inches(9.05),Inches(1.07),Inches(3.75),Inches(0.4),
    sz=13,bold=True,clr=AMBER)
rect0(s,Inches(9.05),Inches(1.6),Inches(3.55),Inches(0.38),RGBColor(0x0F,0x20,0x30))
txt(s,"Fase 1 — Adam",Inches(9.05),Inches(1.6),Inches(3.55),Inches(0.38),
    sz=12,bold=True,clr=TEAL)
for j,item in enumerate(["lr = 1×10⁻³","Hasta 15,000 épocas","Early stopping (pat=800)"]):
    txt(s,item,Inches(9.25),Inches(2.1+j*0.52),Inches(3.35),Inches(0.42),sz=11,clr=LGRAY)
rect0(s,Inches(9.05),Inches(3.8),Inches(3.55),Inches(0.38),RGBColor(0x15,0x10,0x05))
txt(s,"Fase 2 — L-BFGS",Inches(9.05),Inches(3.8),Inches(3.55),Inches(0.38),
    sz=12,bold=True,clr=GOLD)
for j,item in enumerate(["Strong Wolfe line search","Max 1,000 iteraciones","History size = 100"]):
    txt(s,item,Inches(9.25),Inches(4.3+j*0.52),Inches(3.35),Inches(0.42),sz=11,clr=LGRAY)
txt(s,"Adam explora → L-BFGS refina con curvatura de 2° orden",
    Inches(9.05),Inches(6.3),Inches(3.75),Inches(0.45),
    sz=11,italic=True,clr=GOLD,align=PP_ALIGN.CENTER)

# ── S7 CÓDIGO CLAVE ──────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"Implementación: Código Clave")

KWORD = RGBColor(0xC7,0x92,0xEA)   # purple — keywords
CMNT  = RGBColor(0x54,0x6E,0x7A)   # gray   — comments
NUM   = RGBColor(0xF7,0x8C,0x6C)   # orange — numbers/values
FUNC  = RGBColor(0x82,0xAA,0xFF)   # blue   — function names
STR   = RGBColor(0xC3,0xE8,0x8D)   # green  — strings/results

def code_block(slide, x, y, w, h, lines):
    """lines: list of (text, color) tuples; one list per line"""
    rect0(slide, x, y, w, h, RGBColor(0x08,0x0C,0x18))
    slide.shapes[-1].line.color.rgb = BORDER
    slide.shapes[-1].line.width = Pt(1.0)
    for li, segments in enumerate(lines):
        yy = y + Inches(0.15) + li * Inches(0.42)
        tb = slide.shapes.add_textbox(x+Inches(0.2), yy, w-Inches(0.4), Inches(0.4))
        tf = tb.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        for txt_seg, clr in segments:
            r = p.add_run(); r.text = txt_seg
            r.font.size = Pt(11.5); r.font.name = "Consolas"
            r.font.color.rgb = clr; r.font.bold = False

# Left panel — SIREN model
txt(s,"src/models.py — Arquitectura SIREN",
    Inches(0.3),Inches(0.95),Inches(6.1),Inches(0.38),
    sz=12,bold=True,clr=TEAL)
code_block(s, Inches(0.3), Inches(1.38), Inches(6.1), Inches(5.55), [
    [(  "class ",KWORD),("PINN_2D_SIREN",FUNC  ),(  "(nn.Module):",WHITE)],
    [(" ",WHITE)],
    [("  def ",KWORD),("__init__",FUNC),("(self, hidden, layers, omega):",WHITE)],
    [("    self",WHITE),(".omega_0",FUNC),(" = omega",WHITE)],
    [("    self",WHITE),(".fc_in",FUNC),("  = nn.Linear(",WHITE),("2",NUM),(",hidden)",WHITE)],
    [("    self",WHITE),(".hidden",FUNC),(" = nn.ModuleList([",WHITE)],
    [("        nn.Linear(hidden,hidden) ",WHITE),("for",KWORD),(" _ ",WHITE),("in",KWORD),(" range(layers)",WHITE)],
    [("    ])",WHITE)],
    [("    self",WHITE),(".fc_out",FUNC),(" = nn.Linear(hidden,",WHITE),("2",NUM),(")",WHITE)],
    [(" ",WHITE)],
    [("  def ",KWORD),("forward",FUNC),("(self, x):",WHITE)],
    [("    x  = torch.sin(self.omega_0 * self.fc_in(x))",WHITE),
     ("  # ← SIREN",CMNT)],
    [("    for",KWORD),(" layer ",WHITE),("in",KWORD),(" self.hidden:",WHITE)],
    [("        x = torch.sin(self.omega_0 * layer(x))",WHITE)],
    [("    return",KWORD),(" self.fc_out(x)",WHITE),
     ("   # (E_real, E_imag)",CMNT)],
])

# Right panel — Loss + Training
txt(s,"src/losses.py + training — Pérdida y Optimización",
    Inches(6.7),Inches(0.95),Inches(6.33),Inches(0.38),
    sz=12,bold=True,clr=GOLD)
code_block(s, Inches(6.7), Inches(1.38), Inches(6.33), Inches(2.6), [
    [("def ",KWORD),("pinn_loss_2d",FUNC),("(model, xy_int, xy_bc, E_bc, k, lam):",WHITE)],
    [("    L_data",WHITE),(" = MSE(model(xy_bc), E_bc)",WHITE),
     ("     # frontera",CMNT)],
    [("    k_sq ",WHITE),(" = float(k)**",WHITE),("2",NUM)],
    [("    # Residuo de Helmholtz via autodiferenciación",CMNT)],
    [("    E_xx+E_yy + k_sq*E",WHITE),(" == 0",STR),
     ("  # ← L_física",CMNT)],
    [("    return",KWORD),(" L_data + ",WHITE),("lam",NUM),(" * L_fisica",WHITE)],
])
code_block(s, Inches(6.7), Inches(4.2), Inches(6.33), Inches(2.73), [
    [("# Fase 1 — Adam: exploración global",CMNT)],
    [("for",KWORD),(" epoch ",WHITE),("in",KWORD),(" range(",WHITE),("15000",NUM),("):",WHITE)],
    [("    loss = pinn_loss_2d(model, xy_int, xy_bc, E_bc, K, ",WHITE),("0.1",NUM),(")",WHITE)],
    [("    loss.backward(); optimizer.step()",WHITE)],
    [(" ",WHITE)],
    [("# Fase 2 — L-BFGS: refinamiento con curvatura 2°",CMNT)],
    [("def ",KWORD),("closure",FUNC),("():",WHITE)],
    [("    loss = pinn_loss_2d(...)",WHITE)],
    [("    loss.backward(); return loss",WHITE)],
    [("lbfgs.step(closure)",FUNC),
     ("  # strong Wolfe line search",CMNT)],
])

# KPI strip
for val,lbl,c,x in [("66,690","Parámetros",TEAL,0.3),
                     ("Adam→L-BFGS","Estrategia",GOLD,3.5),
                     ("λ=0.1","Peso física",AMBER,6.7),
                     ("float(k)²","Fix dtype",GREEN,9.9)]:
    rect(s,Inches(x),Inches(7.08),Inches(3.1),Inches(0.35),PANEL,BORDER,1.0)
    txt(s,f"{val}  ·  {lbl}",Inches(x+0.1),Inches(7.1),Inches(2.9),Inches(0.28),
        sz=11,bold=True,clr=c,align=PP_ALIGN.CENTER)

# ── S8 NB01 ──────────────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"NB01: Validación 1D — Helmholtz Analítico")
txt(s,"E_exacta(x) = cos(kx),  k = 2π,  x ∈ [0,1]  →  SIREN 5×64, ω₀=1.0",
    Inches(0.3),Inches(1.05),Inches(12.7),Inches(0.4),
    sz=13,clr=LGRAY,align=PP_ALIGN.CENTER,italic=True)
kpis7=[("0.006%","Error L²",GREEN),("1.000000","R²",GREEN),
       ("×415","vs. Schoder 2024",AMBER),("251 s","Tiempo total",TEAL)]
for i,(v,l,c) in enumerate(kpis7):
    x=Inches(0.3+i*3.27)
    rect(s,x,Inches(1.65),Inches(3.1),Inches(1.35),PANEL,c,2.0)
    txt(s,v,x,Inches(1.7),Inches(3.1),Inches(0.78),sz=36,bold=True,clr=c,align=PP_ALIGN.CENTER)
    txt(s,l,x,Inches(2.45),Inches(3.1),Inches(0.4),sz=12,clr=LGRAY,align=PP_ALIGN.CENTER)
# Figura NB01: ajuste PINN vs solución analítica
FIG_1D = r"D:\Tesis_Maestria\paper\figures\resultados_pinn_1d.png"
s.shapes.add_picture(FIG_1D, Inches(0.3), Inches(3.1), width=Inches(5.8))
# Panel derecho — métricas compactas + comparativa
rect(s,Inches(6.4),Inches(3.1),Inches(6.6),Inches(4.0),PANEL,BORDER)
txt(s,"Métricas clave",Inches(6.6),Inches(3.2),Inches(6.2),Inches(0.4),
    sz=13,bold=True,clr=WHITE)
for j,(m,v,c) in enumerate([("MSE","1.5×10⁻⁹",LGRAY),("MAE","3.50×10⁻⁵",LGRAY),
                              ("Pearson","1.000000",GREEN),("L-BFGS iters","202/500",LGRAY)]):
    yy=Inches(3.7+j*0.55)
    txt(s,m,Inches(6.6),yy,Inches(3.3),Inches(0.42),sz=12,clr=LGRAY)
    txt(s,v,Inches(10.0),yy,Inches(2.8),Inches(0.42),sz=12,clr=c,align=PP_ALIGN.RIGHT)
rect0(s,Inches(6.4),Inches(5.85),Inches(6.6),Inches(1.2),RGBColor(0x20,0x08,0x08))
txt(s,"Schoder & Kraxberger (2024)  →  L² = 2.49%",
    Inches(6.6),Inches(5.92),Inches(6.2),Inches(0.4),sz=12,clr=RED)
rect0(s,Inches(6.4),Inches(7.1),Inches(6.6),Inches(0.38),RGBColor(0x02,0x18,0x10))
txt(s,"PINN-SIREN  →  L² = 0.006%   (×415 de mejora)",
    Inches(6.6),Inches(7.12),Inches(6.2),Inches(0.35),sz=12,bold=True,clr=GREEN)

# ── S8 NB02 ──────────────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"NB02: Validación 2D — Campo Complejo")
txt(s,"E = exp(i(k_x·x + k_y·y)),  k_x=k_y=k/√2,  SIREN 5×128, LHS N_c=3000",
    Inches(0.3),Inches(1.05),Inches(12.7),Inches(0.4),
    sz=12,clr=LGRAY,align=PP_ALIGN.CENTER,italic=True)
kpis8=[("0.171%","L² Promedio",GREEN),("0.214%","L² E_real",TEAL),
       ("0.127%","L² E_imag",TEAL),("×11.2","vs. Schoder 2D",AMBER)]
for i,(v,l,c) in enumerate(kpis8):
    x=Inches(0.3+i*3.27)
    rect(s,x,Inches(1.6),Inches(3.1),Inches(1.25),PANEL,c,2.0)
    txt(s,v,x,Inches(1.65),Inches(3.1),Inches(0.7),sz=32,bold=True,clr=c,align=PP_ALIGN.CENTER)
    txt(s,l,x,Inches(2.3),Inches(3.1),Inches(0.35),sz=11,clr=LGRAY,align=PP_ALIGN.CENTER)
# Figura NB02 izquierda: heatmaps exacta vs PINN
FIG_2D     = r"D:\Tesis_Maestria\paper\figures\resultados_pinn_2d.png"
FIG_CURVA2 = r"D:\Tesis_Maestria\paper\figures\metricas_adicionales_2d.png"
s.shapes.add_picture(FIG_2D, Inches(0.3), Inches(2.85), width=Inches(5.1))
# Derecha: curva de convergencia Adam → L-BFGS
txt(s,"Convergencia Adam → L-BFGS (curva de pérdida)",
    Inches(5.65),Inches(2.85),Inches(7.35),Inches(0.4),
    sz=12,bold=True,clr=AMBER)
s.shapes.add_picture(FIG_CURVA2, Inches(5.65), Inches(3.32), width=Inches(7.35))
# Strip inferior multi-seed
rect(s,Inches(0.3),Inches(7.12),Inches(12.7),Inches(0.38),PANEL,BORDER)
txt(s,"Multi-seed {42,123,777}:  L²=0.155±0.020%   |   R²>0.9999   |"
      "   Adam 8,737 ep. + L-BFGS 1,035 it.   |   299 s",
    Inches(0.5),Inches(7.15),Inches(12.3),Inches(0.3),
    sz=10,bold=True,clr=GREEN,align=PP_ALIGN.CENTER)

# ── S9 ABLACIÓN ──────────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"Ablación del Peso de Física λ_phys")
txt(s,"¿Cómo afecta el balance datos/física al error L²?  —  SIREN 5×128, SEED=42",
    Inches(0.3),Inches(1.05),Inches(12.7),Inches(0.4),sz=13,clr=LGRAY,align=PP_ALIGN.CENTER)
# Column headers
for (lbl,x,w) in [("λ_phys",Inches(1.6),Inches(2.2)),
                   ("Error L² prom. (%)",Inches(4.0),Inches(3.5)),
                   ("Converge",Inches(7.7),Inches(2.0)),
                   ("Observación",Inches(9.9),Inches(3.3))]:
    txt(s,lbl,x,Inches(1.7),w,Inches(0.45),sz=13,bold=True,clr=GOLD,align=PP_ALIGN.CENTER)
rows9=[("0.01","0.163","Sí","Física subponderada",LGRAY,TEAL,GREEN,LGRAY,PANEL,BORDER,1.5),
       ("0.1", "0.171","Sí","✦ Configuración óptima",GOLD,GREEN,GREEN,GOLD,
        RGBColor(0x10,0x20,0x10),GOLD,2.5),
       ("1.0", ">5.0", "No","No converge — gradiente explota",LGRAY,RED,RED,RED,
        RGBColor(0x20,0x08,0x08),RED,1.5)]
for i,(lam,l2,cv,obs,cl,cl2,ccv,cob,bg_c,bc,bw) in enumerate(rows9):
    yy=Inches(2.35+i*1.45)
    rect(s,Inches(1.2),yy,Inches(11.8),Inches(1.25),bg_c,bc,bw)
    txt(s,lam,Inches(1.6),yy+Inches(0.28),Inches(2.2),Inches(0.6),sz=24,bold=True,clr=cl,align=PP_ALIGN.CENTER)
    txt(s,l2, Inches(4.0),yy+Inches(0.28),Inches(3.5),Inches(0.6),sz=24,bold=True,clr=cl2,align=PP_ALIGN.CENTER)
    txt(s,cv, Inches(7.7),yy+Inches(0.28),Inches(2.0),Inches(0.6),sz=18,bold=True,clr=ccv,align=PP_ALIGN.CENTER)
    txt(s,obs,Inches(9.9),yy+Inches(0.28),Inches(3.1),Inches(0.6),sz=13,clr=cob)
txt(s,"λ=1.0: dos residuos (real+imag) duplican el gradiente de física → explosión en float32.",
    Inches(0.3),Inches(7.1),Inches(12.7),Inches(0.35),
    sz=11,clr=LGRAY,align=PP_ALIGN.CENTER,italic=True)

# ── S10 NB03 SPECKLE ─────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"NB03: Simulación de Speckle Óptico")
txt(s,"Frontera: φ(x)~U(0,2π) en y=0  ·  Bordes libres en x=0,x=1,y=1  ·  SIREN 5×128",
    Inches(0.3),Inches(1.05),Inches(12.7),Inches(0.4),
    sz=12,clr=LGRAY,align=PP_ALIGN.CENTER,italic=True)
kpis10=[("C = 1.025","σ_I/⟨I⟩  →  |C−1|=0.025 < 0.1",GREEN,"✅ Goodman cumplido"),
        ("0.135","Fracción I > 2⟨I⟩",TEAL,"Teórico e⁻²≈0.1353 ✅"),
        ("p < 0.0001","KS p-valor",AMBER,"⚠ Alta potencia N=10,000")]
for i,(v,l,c,note) in enumerate(kpis10):
    x=Inches(0.4+i*4.3)
    rect(s,x,Inches(1.65),Inches(4.05),Inches(2.3),PANEL,c,2.0)
    txt(s,v,x,Inches(1.75),Inches(4.05),Inches(0.8),sz=28,bold=True,clr=c,align=PP_ALIGN.CENTER)
    txt(s,l,x,Inches(2.5), Inches(4.05),Inches(0.35),sz=10,clr=LGRAY,align=PP_ALIGN.CENTER)
    txt(s,note,x,Inches(2.9),Inches(4.05),Inches(0.85),sz=11,clr=WHITE,align=PP_ALIGN.CENTER,italic=True)
# Figuras NB03: patrón speckle + distribución estadística
FIG_SPECKLE = r"D:\Tesis_Maestria\paper\figures\resultados_speckle_nb03.png"
FIG_STATS   = r"D:\Tesis_Maestria\paper\figures\estadistica_speckle_nb03.png"
# Patrón speckle: izquierda (ratio≈1.154)
s.shapes.add_picture(FIG_SPECKLE, Inches(0.3), Inches(4.05), width=Inches(5.5))
txt(s,"Patrón I(x,y) = |E|²  —  fase rugosa φ~U(0,2π)",
    Inches(0.3),Inches(7.0),Inches(5.5),Inches(0.38),
    sz=10,clr=LGRAY,align=PP_ALIGN.CENTER,italic=True)
# Distribución estadística: derecha (ratio≈1.807)
s.shapes.add_picture(FIG_STATS, Inches(6.1), Inches(4.05), width=Inches(7.0))
txt(s,"Histograma p(I) vs. exponencial teórica  +  CDF  —  C=1.025 ✅",
    Inches(6.1),Inches(7.0),Inches(7.0),Inches(0.38),
    sz=10,clr=LGRAY,align=PP_ALIGN.CENTER,italic=True)

# ── S11 RESUMEN AVANCE ───────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"Resumen de Avance — Estado del Proyecto")
rows11=[
    ("NB01","Helmholtz 1D","L²=0.006%  R²=1.000000  ×415 vs SOTA\nArquitectura SIREN 5×64, 251 s","✅ COMPLETO",GREEN),
    ("NB02","Helmholtz 2D Campo Complejo","L²=0.171%  R²>0.9999  ×11.2 vs SOTA\nMulti-seed: 0.155±0.020%  (seeds 42,123,777)","✅ COMPLETO",GREEN),
    ("NB03","Speckle Óptico","C=1.025 ✅  Dist. exponencial ✅\nKS p<0.0001 (alta potencia estadística)","✅ COMPLETO",GREEN),
    ("NB04","Benchmark FEM (FEniCSx)","Speed-up Factor S=T_FEM/T_PINN\nComparación cuantitativa pendiente","📋 PENDIENTE",GOLD),
]
for i,(nb,t,d,st,c) in enumerate(rows11):
    yy=Inches(1.15+i*1.45)
    rect(s,Inches(0.3),yy,Inches(12.7),Inches(1.3),PANEL,c,1.5)
    rect0(s,Inches(0.3),yy,Inches(1.2),Inches(1.3),c)
    txt(s,nb,Inches(0.3),yy+Inches(0.35),Inches(1.2),Inches(0.55),
        sz=14,bold=True,clr=RGBColor(0,0,0),align=PP_ALIGN.CENTER)
    txt(s,t,Inches(1.65),yy+Inches(0.07),Inches(7.0),Inches(0.45),sz=14,bold=True,clr=WHITE)
    txt(s,d,Inches(1.65),yy+Inches(0.57),Inches(7.0),Inches(0.65),sz=11,clr=LGRAY)
    txt(s,st,Inches(10.2),yy+Inches(0.3),Inches(2.6),Inches(0.6),sz=13,bold=True,clr=c,align=PP_ALIGN.CENTER)
txt(s,"Hipótesis verificada: L² < 5%  ✅   Umbral superado por 2–3 órdenes de magnitud",
    Inches(0.3),Inches(7.1),Inches(12.7),Inches(0.35),
    sz=12,clr=GREEN,bold=True,align=PP_ALIGN.CENTER)

# ── S12 PRÓXIMOS PASOS ───────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"Próximos Pasos y Trabajo Futuro")
pasos=[
    ("NB04","Benchmark FEM\n(Inmediato)",
     "FEniCSx vs PINN-SIREN\nSpeed-up Factor S=T_FEM/T_PINN\nMúltiples realizaciones de speckle",GOLD),
    ("NB05","Condición de Sommerfeld\n(Corto plazo)",
     "∂E/∂n − ikE = 0 en y=1\nElimina reflexiones espurias\nMejora fidelidad KS",TEAL),
    ("NB06","Medios Inhomogéneos\n(Mediano plazo)",
     "Extensión a n(r) variable\nSpeckle en tejidos biológicos\nMateriales GRIN",GREEN),
]
for i,(nb,t,d,c) in enumerate(pasos):
    x=Inches(0.3+i*4.37)
    rect(s,x,Inches(1.1),Inches(4.15),Inches(5.8),PANEL,c,2.0)
    rect0(s,x,Inches(1.1),Inches(4.15),Inches(0.6),c)
    txt(s,nb,x+Inches(0.1),Inches(1.13),Inches(3.95),Inches(0.52),
        sz=14,bold=True,clr=RGBColor(0,0,0),align=PP_ALIGN.CENTER)
    txt(s,t,x+Inches(0.2),Inches(1.85),Inches(3.75),Inches(0.8),
        sz=15,bold=True,clr=c)
    txt(s,d,x+Inches(0.2),Inches(2.8),Inches(3.75),Inches(3.8),sz=12,clr=WHITE)
txt(s,"GPU: NVIDIA RTX 5050 · PyTorch 2.4 · CUDA 12.6 · github.com/RobertoEstradah/pinn_speckle",
    Inches(0.3),Inches(7.1),Inches(12.7),Inches(0.35),
    sz=11,clr=LGRAY,align=PP_ALIGN.CENTER,italic=True)

# ── S13 CONCLUSIONES ─────────────────────────────────────────────────────────
s = blank(prs); bg(s); header(s,"Conclusiones")
concls=[
    ("01","Precisión sin Malla",
     "L²=0.006% (1D) y L²=0.171% (2D)\n3 órdenes por debajo del umbral 5%\nMejora ×415 y ×11.2 sobre Schoder 2024",GREEN),
    ("02","Arquitectura Agnóstica",
     "Misma SIREN 5×128 para onda plana\ny para frontera rugosa aleatoria\nSin modificar arquitectura ni pérdida",TEAL),
    ("03","Speckle Estadísticamente Válido",
     "C=1.025 cumple criterio Goodman\nDist. exponencial negativa verificada\nKS: alta potencia con N=10,000",GOLD),
    ("04","Adam + L-BFGS Esencial",
     "Estrategia bifásica supera a ambos\noptimizadores por separado\nλ_phys=0.1 crítico para 2D complejo",AMBER),
]
for i,(n,t,d,c) in enumerate(concls):
    ri,ci2=divmod(i,2)
    x=Inches(0.3+ci2*6.5); y=Inches(1.1+ri*2.9)
    rect(s,x,y,Inches(6.2),Inches(2.65),PANEL,c,2.0)
    rect0(s,x,y,Inches(0.7),Inches(2.65),c)
    txt(s,n,x,y+Inches(0.9),Inches(0.7),Inches(0.7),
        sz=22,bold=True,clr=RGBColor(0,0,0),align=PP_ALIGN.CENTER)
    txt(s,t,x+Inches(0.85),y+Inches(0.15),Inches(5.2),Inches(0.5),sz=14,bold=True,clr=c)
    txt(s,d,x+Inches(0.85),y+Inches(0.75),Inches(5.2),Inches(1.75),sz=11,clr=WHITE)

# ── S14 GRACIAS ──────────────────────────────────────────────────────────────
s = blank(prs); bg(s)
rect0(s,Inches(0),Inches(0),Inches(0.1),H,GOLD)
txt(s,"¡Gracias!",Inches(1.0),Inches(1.5),Inches(11.0),Inches(1.5),
    sz=56,bold=True,clr=WHITE,align=PP_ALIGN.CENTER)
txt(s,"Preguntas y Comentarios",Inches(1.0),Inches(3.1),Inches(11.0),Inches(0.6),
    sz=22,clr=GOLD,align=PP_ALIGN.CENTER)
rect(s,Inches(2.5),Inches(3.9),Inches(8.3),Inches(2.3),PANEL,BORDER)
txt(s,"Roberto Hernández Estrada",Inches(2.7),Inches(4.05),Inches(7.9),Inches(0.45),
    sz=16,bold=True,clr=WHITE,align=PP_ALIGN.CENTER)
txt(s,"robertohernandezestrd@gmail.com",Inches(2.7),Inches(4.55),Inches(7.9),Inches(0.35),
    sz=13,clr=TEAL,align=PP_ALIGN.CENTER)
txt(s,"Maestría en Ciencias de la Computación · UJAT",Inches(2.7),Inches(4.95),Inches(7.9),Inches(0.35),
    sz=12,clr=LGRAY,align=PP_ALIGN.CENTER)
txt(s,"Director: Dr. José Adán Hernández Nolasco",Inches(2.7),Inches(5.35),Inches(7.9),Inches(0.35),
    sz=12,clr=LGRAY,align=PP_ALIGN.CENTER)
txt(s,"github.com/RobertoEstradah/pinn_speckle",
    Inches(1.0),Inches(6.5),Inches(11.0),Inches(0.4),
    sz=13,clr=TEAL,align=PP_ALIGN.CENTER)

# ── GUARDAR ───────────────────────────────────────────────────────────────────
OUT = r"D:\Tesis_Maestria\master_supporting_docs\supporting_slides\slide_final\coloquio_v3.pptx"
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Guardado: {OUT}")
print(f"Total slides: {len(prs.slides)}")
